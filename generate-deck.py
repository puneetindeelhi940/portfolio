#!/usr/bin/env python3
"""
FeedIQ & The Cargill Terminal (T-C-T) — McKinsey-style case study deck (Double Diamond framework)
Author: Puneet Arora
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Colors ──
BG_DARK    = RGBColor(0x0A, 0x0F, 0x0D)
BG_CARD    = RGBColor(0x14, 0x1E, 0x1A)
BG_CARD2   = RGBColor(0x1A, 0x28, 0x22)
ACCENT     = RGBColor(0x00, 0xC9, 0x78)
ACCENT_DIM = RGBColor(0x00, 0xA8, 0x63)
WHITE      = RGBColor(0xE8, 0xF0, 0xEC)
DIM        = RGBColor(0x8F, 0xA8, 0x9C)
FAINT      = RGBColor(0x5A, 0x72, 0x68)
RULE       = RGBColor(0x24, 0x35, 0x30)
WARN       = RGBColor(0xFF, 0x9F, 0x43)
DANGER     = RGBColor(0xFF, 0x6B, 0x6B)
INFO       = RGBColor(0x54, 0xA0, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]


# ══════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.02
    return shape

def add_text(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_label(slide, left, top, text, color=ACCENT):
    return add_text(slide, left, top, Inches(5), Inches(0.35), text,
                    font_size=10, color=color, bold=True, font_name='Consolas')

def add_title_text(slide, left, top, text, width=Inches(10), font_size=32):
    line_count = text.count('\n') + 1
    height = Inches(0.55 * line_count + 0.1)
    return add_text(slide, left, top, width, height, text,
                    font_size=font_size, color=WHITE, bold=True)

def add_body_text(slide, left, top, text, width=Inches(8), font_size=14, color=DIM):
    line_count = text.count('\n') + 1
    height = max(Inches(0.4), Inches(0.22 * line_count + 0.1))
    return add_text(slide, left, top, width, height, text,
                    font_size=font_size, color=color)

def add_stat_card(slide, left, top, width, stat, label, desc="", stat_color=ACCENT):
    card = add_rect(slide, left, top, width, Inches(1.8), BG_CARD, RULE)
    add_text(slide, left + Inches(0.25), top + Inches(0.12), width - Inches(0.5), Inches(0.5),
             stat, font_size=32, color=stat_color, bold=True, font_name='Consolas')
    add_text(slide, left + Inches(0.25), top + Inches(0.65), width - Inches(0.5), Inches(0.25),
             label, font_size=8, color=FAINT, bold=True, font_name='Consolas')
    if desc:
        add_text(slide, left + Inches(0.25), top + Inches(0.95), width - Inches(0.5), Inches(0.75),
                 desc, font_size=9, color=DIM)

def add_kpi_row(slide, left, top, items, total_width=Inches(11.5)):
    card_w = total_width / len(items)
    for i, (val, label) in enumerate(items):
        x = left + card_w * i
        add_rect(slide, x, top, card_w - Inches(0.08), Inches(0.85), BG_CARD2, RULE)
        add_text(slide, x + Inches(0.15), top + Inches(0.08), card_w - Inches(0.3), Inches(0.4),
                 val, font_size=20, color=WHITE, bold=True, font_name='Consolas', alignment=PP_ALIGN.CENTER)
        add_text(slide, x + Inches(0.15), top + Inches(0.48), card_w - Inches(0.3), Inches(0.25),
                 label, font_size=8, color=FAINT, bold=True, font_name='Consolas', alignment=PP_ALIGN.CENTER)

def add_bullet_card(slide, left, top, width, height, title, bullets, title_color=WHITE):
    card = add_rect(slide, left, top, width, height, BG_CARD, RULE)
    add_text(slide, left + Inches(0.25), top + Inches(0.15), width - Inches(0.5), Inches(0.3),
             title, font_size=12, color=title_color, bold=True)
    y = top + Inches(0.5)
    for b in bullets:
        add_text(slide, left + Inches(0.25), y, width - Inches(0.5), Inches(0.45),
                 f"  {b}", font_size=10, color=DIM)
        y += Inches(0.42)
    return card

def add_phase_label(slide, left, top, phase_text, diamond_num=1):
    color = ACCENT if diamond_num == 1 else INFO
    add_rect(slide, left, top, Inches(1.6), Inches(0.3), BG_CARD2, color)
    add_text(slide, left + Inches(0.1), top + Inches(0.02), Inches(1.4), Inches(0.25),
             phase_text, font_size=9, color=color, bold=True, font_name='Consolas', alignment=PP_ALIGN.CENTER)

def add_divider(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RULE
    shape.line.fill.background()

def add_attribution(slide, text="Puneet Arora  |  Product Design & AI Strategy"):
    add_text(slide, Inches(0.8), SLIDE_H - Inches(0.55), Inches(5), Inches(0.3),
             text, font_size=8, color=FAINT, font_name='Consolas')

def add_page_num(slide, num, total=30):
    add_text(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.55), Inches(1), Inches(0.3),
             f"{num} / {total}", font_size=8, color=FAINT, font_name='Consolas', alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)

# Accent bar at top
add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)

add_label(sl, Inches(0.8), Inches(1.5), "—  PRODUCT CASE STUDY  |  JUNE 2026")
add_text(sl, Inches(0.8), Inches(2.0), Inches(10), Inches(2),
         "FeedIQ  &  The Cargill Terminal", font_size=56, color=WHITE, bold=True)
add_text(sl, Inches(0.8), Inches(3.3), Inches(10), Inches(1),
         "AI-Powered Feed Intelligence + Enterprise\nDecision Intelligence Platform", font_size=26, color=ACCENT, bold=False)
add_text(sl, Inches(0.8), Inches(4.6), Inches(10), Inches(0.8),
         "Two products, one vision: transforming Cargill's data into AI-native decisions.\n"
         "FeedIQ tackles feed formulation & mycotoxin risk. The Cargill Terminal (T-C-T)\n"
         "scales to enterprise-wide Decision Intelligence — Executive, Procurement & Supply Chain.", font_size=14, color=DIM)

add_divider(sl, Inches(0.8), Inches(6.0), Inches(11.5))
add_text(sl, Inches(0.8), Inches(6.15), Inches(4), Inches(0.3),
         "PUNEET ARORA", font_size=12, color=WHITE, bold=True, font_name='Consolas')
add_text(sl, Inches(0.8), Inches(6.45), Inches(6), Inches(0.3),
         "Principal Product Designer  &  AI Design Leader  |  20 years  |  Dell · Boeing · Software AG · HP R&D",
         font_size=10, color=FAINT, font_name='Consolas')

add_text(sl, Inches(8), Inches(6.15), Inches(4.5), Inches(0.3),
         "FRAMEWORK: DOUBLE DIAMOND", font_size=10, color=ACCENT, bold=True, font_name='Consolas', alignment=PP_ALIGN.RIGHT)
add_text(sl, Inches(8), Inches(6.45), Inches(4.5), Inches(0.3),
         "Discover → Define → Develop → Deliver",
         font_size=10, color=FAINT, font_name='Consolas', alignment=PP_ALIGN.RIGHT)
add_page_num(sl, 1)


# ══════════════════════════════════════════════════════════════
# SLIDE 2: EXECUTIVE SUMMARY
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_label(sl, Inches(0.8), Inches(0.6), "—  EXECUTIVE SUMMARY")
add_title_text(sl, Inches(0.8), Inches(1.0), "The 60-second brief")

add_rect(sl, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.8), BG_CARD, RULE)

items = [
    ("WHAT", "FeedIQ: AI feed formulation with mycotoxin risk overlay. T-C-T: A Bloomberg Terminal-density enterprise platform with 3 decision engines (Executive, Procurement, Supply Chain), intent-based UX, and AI confidence scoring."),
    ("WHY", "Cargill analyzes 389,926 feed samples/year across 41 countries — insights die in static PDFs. 71% contamination, 8,000 experts departing, $160B revenue with zero AI-native decision tools. The gap between data and action costs billions."),
    ("HOW", "FeedIQ: Conversational AI + constraint optimization + mycotoxin data overlay. T-C-T: 3 deep decision engines with 9 intent modes (Observe → Predict → Execute), AI recommendation cards, supplier intelligence, and shipment tracking."),
    ("WHEN", "Phase 1: Both prototypes DONE. Phase 2 (3-9 mo): Integrate Cargill APIs — mycotoxin DB, procurement data, shipment feeds. Phase 3 (9-18 mo): Production deployment across Cargill's animal nutrition and supply chain operations.")
]

y = Inches(2.2)
for label, desc in items:
    add_text(sl, Inches(1.1), y, Inches(1.2), Inches(0.3),
             label, font_size=10, color=ACCENT, bold=True, font_name='Consolas')
    add_text(sl, Inches(2.5), y, Inches(9.3), Inches(1.15),
             desc, font_size=10, color=DIM)
    y += Inches(1.25)
    if label != "WHEN":
        add_divider(sl, Inches(1.1), y - Inches(0.15), Inches(10.7))

add_attribution(sl)
add_page_num(sl, 2)


# ══════════════════════════════════════════════════════════════
# SLIDE 3: DOUBLE DIAMOND OVERVIEW
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_label(sl, Inches(0.8), Inches(0.6), "—  DESIGN METHODOLOGY")
add_title_text(sl, Inches(0.8), Inches(1.0), "Double Diamond Framework")
add_body_text(sl, Inches(0.8), Inches(1.7),
              "Structured divergent-convergent thinking applied to Cargill's feed intelligence gap.",
              width=Inches(8), font_size=13)

# Diamond 1
add_rect(sl, Inches(0.8), Inches(2.6), Inches(5.5), Inches(4.2), BG_CARD, ACCENT)
add_text(sl, Inches(1.1), Inches(2.75), Inches(5), Inches(0.3),
         "DIAMOND 1: PROBLEM SPACE", font_size=11, color=ACCENT, bold=True, font_name='Consolas')

add_rect(sl, Inches(1.1), Inches(3.2), Inches(2.4), Inches(3.3), BG_CARD2, RULE)
add_text(sl, Inches(1.3), Inches(3.3), Inches(2), Inches(0.25),
         "DISCOVER", font_size=10, color=ACCENT, bold=True, font_name='Consolas')
add_text(sl, Inches(1.3), Inches(3.6), Inches(2), Inches(0.2),
         "Divergent", font_size=9, color=FAINT, font_name='Consolas')
add_text(sl, Inches(1.3), Inches(3.9), Inches(2), Inches(2.4),
         "• Cargill business analysis\n• Industry landscape (ABCD)\n• Stakeholder pain points\n• Mycotoxin data deep dive\n• Competitor AI initiatives\n• Feed formulation market",
         font_size=10, color=DIM)

add_rect(sl, Inches(3.7), Inches(3.2), Inches(2.4), Inches(3.3), BG_CARD2, RULE)
add_text(sl, Inches(3.9), Inches(3.3), Inches(2), Inches(0.25),
         "DEFINE", font_size=10, color=ACCENT, bold=True, font_name='Consolas')
add_text(sl, Inches(3.9), Inches(3.6), Inches(2), Inches(0.2),
         "Convergent", font_size=9, color=FAINT, font_name='Consolas')
add_text(sl, Inches(3.9), Inches(3.9), Inches(2), Inches(2.4),
         "• Problem statement\n• Key insight articulation\n• Opportunity sizing\n• Competitive gap map\n• Target user definition\n• Success metrics",
         font_size=10, color=DIM)

# Diamond 2
add_rect(sl, Inches(6.6), Inches(2.6), Inches(5.8), Inches(4.2), BG_CARD, INFO)
add_text(sl, Inches(6.9), Inches(2.75), Inches(5), Inches(0.3),
         "DIAMOND 2: SOLUTION SPACE", font_size=11, color=INFO, bold=True, font_name='Consolas')

add_rect(sl, Inches(6.9), Inches(3.2), Inches(2.55), Inches(3.3), BG_CARD2, RULE)
add_text(sl, Inches(7.1), Inches(3.3), Inches(2.2), Inches(0.25),
         "DEVELOP", font_size=10, color=INFO, bold=True, font_name='Consolas')
add_text(sl, Inches(7.1), Inches(3.6), Inches(2.2), Inches(0.2),
         "Divergent", font_size=9, color=FAINT, font_name='Consolas')
add_text(sl, Inches(7.1), Inches(3.9), Inches(2.2), Inches(2.4),
         "• Solution ideation (10 ideas)\n• Feature prioritization\n• AI interaction patterns\n• Dashboard design\n• Risk-aware formulation\n• UX principles",
         font_size=10, color=DIM)

add_rect(sl, Inches(9.65), Inches(3.2), Inches(2.55), Inches(3.3), BG_CARD2, RULE)
add_text(sl, Inches(9.85), Inches(3.3), Inches(2.2), Inches(0.25),
         "DELIVER", font_size=10, color=INFO, bold=True, font_name='Consolas')
add_text(sl, Inches(9.85), Inches(3.6), Inches(2.2), Inches(0.2),
         "Convergent", font_size=9, color=FAINT, font_name='Consolas')
add_text(sl, Inches(9.85), Inches(3.9), Inches(2.2), Inches(2.4),
         "• Working prototype\n• 4 scenario demos\n• Metrics framework\n• Roadmap (3 phases)\n• Benchmark validation\n• Presentation & pitch",
         font_size=10, color=DIM)

add_text(sl, Inches(0.8), Inches(6.95), Inches(11.5), Inches(0.3),
         "Designed by Puneet Arora  |  Applied to Cargill's Animal Nutrition & Feed Intelligence domain",
         font_size=9, color=FAINT, font_name='Consolas', alignment=PP_ALIGN.CENTER)
add_page_num(sl, 3)


# ══════════════════════════════════════════════════════════════
# SLIDE 4: CARGILL AT A GLANCE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DISCOVER", 1)
add_label(sl, Inches(2.6), Inches(0.53), "—  CARGILL CONTEXT")
add_title_text(sl, Inches(0.8), Inches(1.0), "Cargill at a glance: the world's\nlargest private company")

add_kpi_row(sl, Inches(0.8), Inches(2.3), [
    ("$154B", "REVENUE (FY2025)"),
    ("160,000+", "EMPLOYEES"),
    ("70+", "COUNTRIES"),
    ("41", "TESTING COUNTRIES"),
    ("#1", "PRIVATE COMPANY"),
])

biz_segments = [
    ("Agricultural Trading", "Grain, oilseeds, palm oil, cotton — ABCD dominance (50-60% of global trade)"),
    ("Animal Nutrition", "Feed production, precision nutrition tools (Galleon, CNS, CattleView)"),
    ("Food Ingredients", "Starches, sweeteners, oils, cocoa — ingredients for global food brands"),
    ("Protein & Meat", "Beef, poultry, egg processing — one of the world's largest processors"),
    ("Industrial & Energy", "Biofuels, biodiesel (1.5B liters/year in Brazil), commodity hedging"),
    ("Logistics", "Ocean freight, barge transport, port operations (30x ROI on Port Optimizer)"),
]
x = Inches(0.8)
y = Inches(3.5)
card_w = Inches(3.75)
for i, (title, desc) in enumerate(biz_segments):
    cx = x + (i % 3) * (card_w + Inches(0.1))
    cy = y + (i // 3) * Inches(1.55)
    add_rect(sl, cx, cy, card_w, Inches(1.4), BG_CARD, RULE)
    add_text(sl, cx + Inches(0.2), cy + Inches(0.12), card_w - Inches(0.4), Inches(0.25),
             title, font_size=12, color=WHITE, bold=True)
    add_text(sl, cx + Inches(0.2), cy + Inches(0.45), card_w - Inches(0.4), Inches(0.85),
             desc, font_size=10, color=DIM)

add_attribution(sl)
add_page_num(sl, 4)


# ══════════════════════════════════════════════════════════════
# SLIDE 5: CARGILL'S CHALLENGES
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DISCOVER", 1)
add_label(sl, Inches(2.6), Inches(0.53), "—  BUSINESS CHALLENGES")
add_title_text(sl, Inches(0.8), Inches(1.0), "Cargill is facing its toughest\nperiod in decades")

challenges = [
    ("FINANCIAL PRESSURE", DANGER, [
        "Revenue down to $154B (lowest since 2021)",
        "< 1/3 of businesses met earnings goals",
        "5 units consolidated to 3",
        "8,000 jobs cut (5% of workforce)"
    ]),
    ("SUPPLY CHAIN & COMPLIANCE", WARN, [
        "EUDR compliance deadline Dec 2026",
        "Sued by ClientEarth for due diligence gaps",
        "66% of emissions from land-use change",
        "GPS polygon mapping at scale unfinished"
    ]),
    ("FOOD SAFETY & QUALITY", ACCENT, [
        "71% of feed samples contaminated",
        "47% have 3+ simultaneous toxins",
        "34% exceed performance risk thresholds",
        "Feed recall (Nutrena aflatoxin, 2025)"
    ]),
    ("KNOWLEDGE & TALENT", INFO, [
        "8,000 departing = institutional knowledge loss",
        "46% of labs report staff shortages",
        "Fragmented AI tools (no unified platform)",
        "Small farmers lack tech literacy"
    ]),
]

x = Inches(0.8)
y = Inches(2.3)
card_w = Inches(2.75)
for i, (title, color, bullets) in enumerate(challenges):
    cx = x + i * (card_w + Inches(0.1))
    add_rect(sl, cx, y, card_w, Inches(4.3), BG_CARD, color)
    add_text(sl, cx + Inches(0.2), y + Inches(0.15), card_w - Inches(0.4), Inches(0.25),
             title, font_size=9, color=color, bold=True, font_name='Consolas')
    by = y + Inches(0.55)
    for b in bullets:
        add_text(sl, cx + Inches(0.2), by, card_w - Inches(0.4), Inches(0.85),
                 f"• {b}", font_size=9, color=DIM)
        by += Inches(0.75) if len(b) > 35 else Inches(0.65)

add_text(sl, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.3),
         "Source: Cargill FY2025 results, Impact Report 2025, ClientEarth filings, Global Mycotoxin Report 2025",
         font_size=8, color=FAINT, font_name='Consolas')
add_attribution(sl)
add_page_num(sl, 5)


# ══════════════════════════════════════════════════════════════
# SLIDE 6: THE MYCOTOXIN PROBLEM
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DISCOVER", 1)
add_label(sl, Inches(2.6), Inches(0.53), "—  DEEP DIVE: MYCOTOXIN CRISIS")
add_title_text(sl, Inches(0.8), Inches(1.0), "The silent crisis in animal feed")
add_body_text(sl, Inches(0.8), Inches(1.7),
              "Mycotoxins are toxic metabolites produced by fungi in grains. They cause production losses,\n"
              "animal health issues, and can enter the human food chain through meat, milk, and eggs.",
              width=Inches(10), font_size=12)

# Big stat cards
add_stat_card(sl, Inches(0.8), Inches(2.7), Inches(3.6),
              "71%", "CONTAMINATION RATE",
              "of global feed samples test positive for at least one mycotoxin above detection limits.")
add_stat_card(sl, Inches(4.6), Inches(2.7), Inches(3.6),
              "47%", "MULTI-TOXIN PREVALENCE",
              "of samples contain 3 or more mycotoxins simultaneously — compounding risk.", stat_color=WARN)
add_stat_card(sl, Inches(8.4), Inches(2.7), Inches(3.9),
              "34%", "ABOVE RISK THRESHOLD",
              "exceed performance-based thresholds linked to measurable livestock productivity losses.", stat_color=DANGER)

# Toxin breakdown
add_rect(sl, Inches(0.8), Inches(4.7), Inches(11.5), Inches(2.0), BG_CARD, RULE)
add_text(sl, Inches(1.1), Inches(4.85), Inches(4), Inches(0.25),
         "TOP MYCOTOXINS BY GLOBAL RISK", font_size=9, color=FAINT, bold=True, font_name='Consolas')

toxins = [
    ("DON", "53%", "Deoxynivalenol — #1 global risk, above thresholds in 53% of analyses"),
    ("FUM", "38%", "Fumonisins — rising risk, especially in corn from S. America & SE Asia"),
    ("ZEN", "35%", "Zearalenone — affects reproduction in swine and dairy cattle"),
    ("AFLA", "28%", "Aflatoxins — carcinogenic, highest in tropical/humid storage regions"),
    ("T-2", "18%", "T-2/HT-2 toxins — immunosuppressive, elevated in Northern Europe"),
    ("OTA", "12%", "Ochratoxin A — nephrotoxic, found in cereals and coffee"),
]

tx = Inches(1.1)
ty = Inches(5.15)
for i, (name, pct, desc) in enumerate(toxins):
    cx = tx + (i % 3) * Inches(3.75)
    cy = ty + (i // 3) * Inches(0.85)
    add_text(sl, cx, cy, Inches(0.5), Inches(0.25),
             name, font_size=9, color=ACCENT, bold=True, font_name='Consolas')
    add_text(sl, cx + Inches(0.5), cy, Inches(0.5), Inches(0.25),
             pct, font_size=9, color=WHITE, bold=True, font_name='Consolas')
    add_text(sl, cx + Inches(1.0), cy, Inches(2.6), Inches(0.75),
             desc, font_size=8, color=DIM)

add_text(sl, Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.3),
         "Source: Cargill 2025 Global Mycotoxin Report — 389,926 analyses across 41 countries",
         font_size=8, color=FAINT, font_name='Consolas')
add_page_num(sl, 6)


# ══════════════════════════════════════════════════════════════
# SLIDE 7: FEED FORMULATION COMPLEXITY
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DISCOVER", 1)
add_label(sl, Inches(2.6), Inches(0.53), "—  DEEP DIVE: FEED FORMULATION")
add_title_text(sl, Inches(0.8), Inches(1.0), "Feed formulation is a constraint\nsatisfaction nightmare")

add_body_text(sl, Inches(0.8), Inches(1.9),
              "A single feed formulation must balance 15+ nutritional parameters, cost constraints,\n"
              "local ingredient availability, species-specific needs, climate conditions, and now —\n"
              "mycotoxin risk. Current tools are spreadsheet-era software with steep learning curves.",
              width=Inches(10), font_size=12)

# Market stats
add_kpi_row(sl, Inches(0.8), Inches(3.0), [
    ("$150M", "FEED SOFTWARE MARKET (2024)"),
    ("$300M", "PROJECTED BY 2033"),
    ("8.2%", "CAGR 2026-2033"),
    ("12%", "COST REDUCTION (ADISSEO AI BENCHMARK)"),
])

# Pain points
add_rect(sl, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.8), BG_CARD, RULE)
add_text(sl, Inches(1.1), Inches(4.35), Inches(4), Inches(0.25),
         "STAKEHOLDER PAIN POINTS", font_size=9, color=FAINT, bold=True, font_name='Consolas')

personas = [
    ("FEED NUTRITIONIST", "Balances 15+ nutrients across species/stages.\nUses legacy software (Bestmix, Format Solutions).\nTakes hours to iterate on a single formulation.\nMycotoxin risk checked separately, if at all."),
    ("LIVESTOCK PRODUCER", "Receives formulation as a spec sheet.\nNo visibility into risk trade-offs.\nDiscov mycotoxin issues only after animal losses.\nNeeds simple answers, not spreadsheets."),
    ("PROCUREMENT MANAGER", "Sourcing decisions disconnected from risk data.\nPrice-driven, ignoring contamination patterns.\nNo tool showing 'corn from Region X = Y% risk.'\nCompliance burden increasing (EUDR, ESG)."),
]

px = Inches(1.1)
py = Inches(4.7)
for i, (role, pains) in enumerate(personas):
    cx = px + i * Inches(3.75)
    add_text(sl, cx, py, Inches(3.4), Inches(0.25),
             role, font_size=10, color=ACCENT, bold=True, font_name='Consolas')
    add_text(sl, cx, py + Inches(0.35), Inches(3.4), Inches(2),
             pains, font_size=9, color=DIM)

add_text(sl, Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.3),
         "Source: Feed Formulation Software Market Report 2026, Adisseo AI platform launch (March 2026), Cargill VIV Asia 2025",
         font_size=8, color=FAINT, font_name='Consolas')
add_attribution(sl)
add_page_num(sl, 7)


# ══════════════════════════════════════════════════════════════
# SLIDE 8: COMPETITOR LANDSCAPE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DISCOVER", 1)
add_label(sl, Inches(2.6), Inches(0.53), "—  COMPETITIVE LANDSCAPE")
add_title_text(sl, Inches(0.8), Inches(1.0), "Cargill's AI tools vs. the gap")

# Cargill's existing tools
add_rect(sl, Inches(0.8), Inches(2.1), Inches(5.5), Inches(4.8), BG_CARD, RULE)
add_text(sl, Inches(1.1), Inches(2.25), Inches(5), Inches(0.25),
         "CARGILL'S EXISTING AI PORTFOLIO", font_size=9, color=ACCENT, bold=True, font_name='Consolas')

tools = [
    ("CMAX / Port Optimizer", "Shipping logistics — 30x ROI"),
    ("CarVe", "Computer vision for protein yield"),
    ("CattleView", "Drone + AI cattle welfare assessment"),
    ("Agriness", "Cloud platform for swine farm management"),
    ("Birdoo", "3D camera for broiler weighing"),
    ("Galleon", "AI for poultry gut microbiome nutrition"),
    ("Ask Emma", "Gen AI for product development"),
    ("Taste Tinker", "Gen AI for food concept generation"),
]

ty = Inches(2.6)
for name, desc in tools:
    add_text(sl, Inches(1.1), ty, Inches(2.2), Inches(0.25),
             name, font_size=9, color=WHITE, bold=True, font_name='Consolas')
    add_text(sl, Inches(3.3), ty, Inches(2.8), Inches(0.25),
             desc, font_size=9, color=DIM)
    ty += Inches(0.32)

# The gap
add_rect(sl, Inches(6.5), Inches(2.1), Inches(5.8), Inches(4.8), BG_CARD, WARN)
add_text(sl, Inches(6.8), Inches(2.25), Inches(5), Inches(0.25),
         "THE GAP — WHAT DOESN'T EXIST", font_size=9, color=WARN, bold=True, font_name='Consolas')

gaps = [
    ("No unified risk + formulation tool", "Mycotoxin data and feed formulation live in separate systems. Nutritionists check risk manually (if at all)."),
    ("No natural language interface", "Galleon and CNS require expert training. No conversational AI for formulation — still forms and spreadsheets."),
    ("No interactive mycotoxin dashboard", "389K analyses published as a static PDF report. No drill-down, no alerts, no region-specific risk scoring."),
    ("No cost-risk trade-off visibility", "Procurement decides on price alone. No tool shows 'this corn saves $12 but adds 42% fumonisin risk.'"),
    ("No knowledge capture system", "With 8,000 people leaving, decades of formulation expertise walks out the door. No AI to capture and scale it."),
]

gy = Inches(2.6)
for title, desc in gaps:
    add_text(sl, Inches(6.8), gy, Inches(5.2), Inches(0.25),
             f"✕  {title}", font_size=9, color=WHITE, bold=True)
    add_text(sl, Inches(7.1), gy + Inches(0.26), Inches(4.9), Inches(0.55),
             desc, font_size=8, color=DIM)
    gy += Inches(0.82)

add_attribution(sl)
add_page_num(sl, 8)


# ══════════════════════════════════════════════════════════════
# SLIDE 9: ABCD COMPETITOR ANALYSIS
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DISCOVER", 1)
add_label(sl, Inches(2.6), Inches(0.53), "—  INDUSTRY BENCHMARKS")
add_title_text(sl, Inches(0.8), Inches(1.0), "ABCD giants & competitor AI initiatives")

# Competitor table
headers = ["COMPANY", "REVENUE", "AI MATURITY", "FEED INTEL TOOL?", "GAP FEEDIQ FILLS"]
companies = [
    ["Cargill", "$154B", "Strong (fragmented)", "Galleon, CNS (expert-only)", "Unified NL interface + risk overlay"],
    ["ADM", "$85B", "Moderate", "None public", "Full-stack feed intelligence"],
    ["Bunge", "$55B+", "Low (infra-focused)", "None", "AI-first approach to nutrition"],
    ["Louis Dreyfus", "$50B", "Low-Moderate", "None", "Risk-aware formulation"],
    ["COFCO", "$70B+", "Unknown", "None public", "Global mycotoxin dashboard"],
    ["Adisseo", "$2B", "High (AI platform)", "AI feed optimizer (CN)", "Multi-region risk overlay"],
]

add_rect(sl, Inches(0.8), Inches(2.2), Inches(11.5), Inches(3.8), BG_CARD, RULE)

# Headers
hx = Inches(1.0)
hw = [Inches(1.5), Inches(1.2), Inches(2), Inches(2.5), Inches(3.8)]
for i, h in enumerate(headers):
    add_text(sl, hx, Inches(2.35), hw[i], Inches(0.25),
             h, font_size=9, color=FAINT, bold=True, font_name='Consolas')
    hx += hw[i]

add_divider(sl, Inches(1.0), Inches(2.65), Inches(11.0))

# Rows
ry = Inches(2.8)
for row in companies:
    rx = Inches(1.0)
    for i, cell in enumerate(row):
        color = WHITE if i == 0 else (ACCENT if i == 4 else DIM)
        weight = True if i == 0 else False
        add_text(sl, rx, ry, hw[i], Inches(0.45),
                 cell, font_size=8, color=color, bold=weight, font_name='Consolas' if i < 2 else 'Segoe UI')
        rx += hw[i]
    ry += Inches(0.5)

# Benchmark callout
add_rect(sl, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.6), BG_CARD2, ACCENT)
add_text(sl, Inches(1.1), Inches(6.35), Inches(10.5), Inches(0.5),
         "KEY BENCHMARK: Adisseo (March 2026) launched an AI feed optimizer for Chinese livestock — reported 12% cost reduction for major pig farms. "
         "This validates the market. FeedIQ differentiates with global mycotoxin risk integration (Adisseo doesn't have it).",
         font_size=10, color=DIM)

add_attribution(sl, "Analysis by Puneet Arora  |  Sources: Industry reports, company filings, press releases")
add_page_num(sl, 9)


# ══════════════════════════════════════════════════════════════
# SLIDE 10: PROBLEM STATEMENT (DEFINE)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEFINE", 1)
add_label(sl, Inches(2.6), Inches(0.53), "—  PROBLEM STATEMENT")

add_text(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(2),
         "Cargill tests 390,000\nsamples a year.\nThe insight dies in a PDF.",
         font_size=44, color=WHITE, bold=True)

add_rect(sl, Inches(0.8), Inches(4.0), Inches(11.5), Inches(2.2), BG_CARD, ACCENT)
add_text(sl, Inches(1.1), Inches(4.15), Inches(4), Inches(0.25),
         "THE CORE INSIGHT", font_size=10, color=ACCENT, bold=True, font_name='Consolas')
add_text(sl, Inches(1.1), Inches(4.5), Inches(10.5), Inches(1.5),
         "Cargill has the world's largest proprietary mycotoxin dataset (389,926 analyses, 41 countries).\n"
         "Livestock producers have the world's most expensive contamination problem ($1.5B testing market, "
         "unmeasured billions in production losses).\n\n"
         "But there is no interactive tool connecting Cargill's risk intelligence to feed formulation decisions.\n"
         "The data exists. The problem exists. The bridge does not.\n\n"
         "FeedIQ is that bridge.",
         font_size=14, color=DIM)

add_text(sl, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.3),
         "Problem framed by Puneet Arora  |  Validated against Cargill Impact Report 2025, Global Mycotoxin Report 2025, industry interviews",
         font_size=8, color=FAINT, font_name='Consolas')
add_page_num(sl, 10)


# ══════════════════════════════════════════════════════════════
# SLIDE 11: OPPORTUNITY SIZING
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEFINE", 1)
add_label(sl, Inches(2.6), Inches(0.53), "—  OPPORTUNITY SIZING")
add_title_text(sl, Inches(0.8), Inches(1.0), "Two converging markets,\none product opportunity")

# TAM/SAM/SOM
add_stat_card(sl, Inches(0.8), Inches(2.3), Inches(3.6),
              "$1.55B", "TAM — MYCOTOXIN TESTING MARKET",
              "Global mycotoxin testing market (2025), growing at 6.4% CAGR through 2034.")
add_stat_card(sl, Inches(4.6), Inches(2.3), Inches(3.6),
              "$300M", "SAM — FEED FORMULATION SOFTWARE",
              "Feed formulation software market by 2033, growing at 8.2% CAGR.", stat_color=INFO)
add_stat_card(sl, Inches(8.4), Inches(2.3), Inches(3.9),
              "$50-80M", "SOM — RISK-AWARE FORMULATION",
              "Intersection: AI-powered feed tools with integrated mycotoxin risk. No incumbent.", stat_color=WARN)

# Value drivers
add_rect(sl, Inches(0.8), Inches(4.4), Inches(11.5), Inches(2.5), BG_CARD, RULE)
add_text(sl, Inches(1.1), Inches(4.55), Inches(4), Inches(0.25),
         "VALUE DRIVERS FOR CARGILL", font_size=9, color=FAINT, bold=True, font_name='Consolas')

drivers = [
    ("Revenue: premium service layer", "FeedIQ as a value-add for Cargill's animal nutrition customers — justify premium pricing on feed products by bundling risk intelligence."),
    ("Cost avoidance: reduce production losses", "Mycotoxin-related losses cost the global livestock industry billions annually. Even a 5% reduction = massive ROI."),
    ("Retention: knowledge capture", "With 8,000 employees leaving, FeedIQ captures formulation expertise in an AI system that doesn't retire."),
    ("Compliance: EUDR & ESG defense", "Demonstrate responsible sourcing with data-backed decisions. Shift the narrative from 'Cargill is failing on sustainability' to 'Cargill leads with data.'"),
]

dy = Inches(4.9)
for i, (title, desc) in enumerate(drivers):
    cx = Inches(1.1) + (i % 2) * Inches(5.6)
    cy = dy + (i // 2) * Inches(1.15)
    add_text(sl, cx, cy, Inches(5.2), Inches(0.25),
             title, font_size=10, color=ACCENT, bold=True)
    add_text(sl, cx, cy + Inches(0.28), Inches(5.2), Inches(0.8),
             desc, font_size=9, color=DIM)

add_attribution(sl)
add_page_num(sl, 11)


# ══════════════════════════════════════════════════════════════
# SLIDE 12: KEY INSIGHT
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEFINE", 1)
add_label(sl, Inches(2.6), Inches(0.53), "—  KEY DESIGN INSIGHT")

add_text(sl, Inches(0.8), Inches(1.8), Inches(11), Inches(1.5),
         "\"The best AI products don't\nadd intelligence.\nThey remove ignorance.\"",
         font_size=40, color=ACCENT, bold=True)

add_text(sl, Inches(0.8), Inches(3.8), Inches(10), Inches(0.3),
         "— Design principle applied by Puneet Arora to FeedIQ",
         font_size=12, color=FAINT, font_name='Consolas')

add_rect(sl, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.0), BG_CARD, RULE)
add_text(sl, Inches(1.1), Inches(4.65), Inches(10.5), Inches(1.7),
         "Feed nutritionists aren't unintelligent — they're uninformed.\n\n"
         "They formulate feeds without knowing that the corn they specified carries 56% fumonisin risk "
         "this quarter. They optimize cost without seeing that the cheapest sourcing region has the "
         "highest contamination rate.\n\n"
         "FeedIQ doesn't replace the nutritionist's expertise. It removes the information gap that "
         "makes their expertise incomplete. Every formulation is now risk-aware by default — not "
         "because the user remembered to check, but because the system checked for them.",
         font_size=13, color=DIM)

add_attribution(sl)
add_page_num(sl, 12)


# ══════════════════════════════════════════════════════════════
# SLIDE 13: SOLUTION VISION (DEVELOP)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEVELOP", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  SOLUTION VISION")
add_title_text(sl, Inches(0.8), Inches(1.0), "FeedIQ: from data to decision\nin seconds")

add_body_text(sl, Inches(0.8), Inches(1.9),
              "A unified AI platform that connects mycotoxin risk intelligence with feed formulation —\n"
              "making every feed decision safety-aware, cost-optimized, and explainable.",
              width=Inches(10), font_size=13)

# Three pillars
pillars = [
    ("01 — RISK DASHBOARD", "Interactive mycotoxin heatmap",
     "8 regions × 6 toxins, color-coded by severity. Real-time alerts for spikes. "
     "Multi-toxin prevalence gauges. Drill-down by commodity, region, season.",
     ACCENT),
    ("02 — AI FORMULATOR", "Natural language feed optimization",
     "Describe needs in plain English: species, budget, constraints. AI generates "
     "optimized formulation with ingredient breakdown, nutrition profile, and cost analysis.",
     INFO),
    ("03 — TRADE-OFF ENGINE", "Cost-risk transparency",
     "Interactive sliders to explore budget vs. risk trade-offs. AI explains: 'Switching "
     "corn source adds $12/ton but reduces fumonisin risk from 56% to 8%.'",
     WARN),
]

x = Inches(0.8)
y = Inches(2.8)
for i, (title, subtitle, desc, color) in enumerate(pillars):
    cx = x + i * Inches(3.9)
    add_rect(sl, cx, y, Inches(3.7), Inches(3.8), BG_CARD, color)
    add_text(sl, cx + Inches(0.25), y + Inches(0.2), Inches(3.2), Inches(0.25),
             title, font_size=10, color=color, bold=True, font_name='Consolas')
    add_text(sl, cx + Inches(0.25), y + Inches(0.55), Inches(3.2), Inches(0.3),
             subtitle, font_size=16, color=WHITE, bold=True)
    add_text(sl, cx + Inches(0.25), y + Inches(1.0), Inches(3.2), Inches(2.5),
             desc, font_size=12, color=DIM)

add_attribution(sl, "Product concept & design by Puneet Arora")
add_page_num(sl, 13)


# ══════════════════════════════════════════════════════════════
# SLIDE 14: FEATURE — RISK DASHBOARD
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEVELOP", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  FEATURE: RISK DASHBOARD")
add_title_text(sl, Inches(0.8), Inches(1.0), "Global contamination heatmap\nwith real-time alerts")

# Simulated heatmap
add_rect(sl, Inches(0.8), Inches(2.3), Inches(7.5), Inches(4.5), BG_CARD, RULE)
add_text(sl, Inches(1.1), Inches(2.45), Inches(6), Inches(0.25),
         "CONTAMINATION HEATMAP — REGION × TOXIN (% ABOVE THRESHOLD)", font_size=9, color=FAINT, bold=True, font_name='Consolas')

# Simplified heatmap representation
regions = ["N. America", "S. America", "W. Europe", "C. Europe", "SE Asia", "S. Asia", "Africa", "Oceania"]
toxin_names = ["DON", "FUM", "ZEN", "AFLA", "T-2", "OTA"]
values = [
    [48, 22, 30, 12, 18, 8],
    [58, 45, 35, 28, 15, 10],
    [42, 18, 38, 6,  22, 14],
    [52, 24, 44, 8,  28, 12],
    [38, 56, 28, 42, 12, 18],
    [32, 48, 22, 38, 10, 16],
    [28, 42, 18, 52, 8,  14],
    [12, 8,  10, 4,  6,  5],
]

# Headers
for j, t in enumerate(toxin_names):
    add_text(sl, Inches(2.5 + j * 0.9), Inches(2.75), Inches(0.8), Inches(0.2),
             t, font_size=9, color=FAINT, bold=True, font_name='Consolas', alignment=PP_ALIGN.CENTER)

# Data
for i, (region, row) in enumerate(zip(regions, values)):
    ry = Inches(3.05 + i * 0.44)
    add_text(sl, Inches(1.1), ry, Inches(1.3), Inches(0.25),
             region, font_size=9, color=DIM, font_name='Consolas')
    for j, v in enumerate(row):
        if v > 45: c = DANGER
        elif v > 30: c = WARN
        elif v > 15: c = RGBColor(0x80, 0xC9, 0xA0)
        else: c = ACCENT
        cx = Inches(2.5 + j * 0.9)
        cell_bg = BG_CARD2
        add_rect(sl, cx, ry, Inches(0.8), Inches(0.35), cell_bg, RULE)
        add_text(sl, cx, ry + Inches(0.03), Inches(0.8), Inches(0.25),
                 f"{v}%", font_size=10, color=c, bold=True, font_name='Consolas', alignment=PP_ALIGN.CENTER)

# Alerts sidebar
add_rect(sl, Inches(8.5), Inches(2.3), Inches(3.8), Inches(4.5), BG_CARD, RULE)
add_text(sl, Inches(8.8), Inches(2.45), Inches(3.2), Inches(0.25),
         "ACTIVE ALERTS", font_size=9, color=DANGER, bold=True, font_name='Consolas')

alerts = [
    ("DON spike — S. America", "53%+ threshold in corn/wheat", DANGER),
    ("FUM surge — SE Asia", "+18% vs. prior year", DANGER),
    ("Multi-toxin — C. Europe", "52% of samples have 3+ toxins", WARN),
    ("Low risk — Oceania", "All categories below 15%", ACCENT),
]

ay = Inches(2.8)
for title, desc, color in alerts:
    add_rect(sl, Inches(8.7), ay, Inches(3.4), Inches(0.65), BG_CARD2, RULE)
    add_text(sl, Inches(8.9), ay + Inches(0.05), Inches(3), Inches(0.2),
             title, font_size=10, color=color, bold=True)
    add_text(sl, Inches(8.9), ay + Inches(0.28), Inches(3), Inches(0.3),
             desc, font_size=9, color=DIM)
    ay += Inches(0.75)

add_text(sl, Inches(8.7), Inches(5.85), Inches(3.4), Inches(0.8),
         "UX DECISION: Alerts are not notifications — they're actionable recommendations. "
         "Each alert links to a mitigation (swap source, add binder, adjust formulation).",
         font_size=9, color=FAINT)

add_attribution(sl, "Designed by Puneet Arora  |  Data: Cargill 2025 Global Mycotoxin Report")
add_page_num(sl, 14)


# ══════════════════════════════════════════════════════════════
# SLIDE 15: FEATURE — AI FORMULATOR
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEVELOP", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  FEATURE: AI FORMULATOR")
add_title_text(sl, Inches(0.8), Inches(1.0), "Conversational AI for feed\nformulation")

# Chat mockup
add_rect(sl, Inches(0.8), Inches(2.3), Inches(7), Inches(4.5), BG_CARD, RULE)
add_text(sl, Inches(1.1), Inches(2.45), Inches(6), Inches(0.25),
         "FEEDIQ AI FORMULATOR", font_size=9, color=FAINT, bold=True, font_name='Consolas')

# User message
add_rect(sl, Inches(2.5), Inches(2.85), Inches(5.0), Inches(0.7), ACCENT)
add_text(sl, Inches(2.7), Inches(2.9), Inches(4.6), Inches(0.55),
         "Optimize a broiler feed for tropical climate,\nbudget under $280/ton, minimize soy content",
         font_size=11, color=BG_DARK, bold=False)

# AI response
add_rect(sl, Inches(1.1), Inches(3.7), Inches(5.5), Inches(2.8), BG_CARD2, RULE)
add_text(sl, Inches(1.3), Inches(3.78), Inches(1.5), Inches(0.2),
         "FEEDIQ AI", font_size=8, color=ACCENT, bold=True, font_name='Consolas')
add_text(sl, Inches(1.3), Inches(4.0), Inches(5.1), Inches(2.4),
         "Analyzing: Broiler (grower), Tropical, <$280/ton, Low soy\n\n"
         "Generated formulation at $268/ton — 4.3% under budget.\n\n"
         "Key decisions:\n"
         "• Replaced 30% of soy with rice bran + DDGS\n"
         "• Added palm oil for energy in heat stress\n"
         "• Flagged: SE Asian corn = 56% fumonisin risk\n"
         "• Included 0.5% mycotoxin binder as safety net",
         font_size=10, color=DIM)

# Why conversational
add_rect(sl, Inches(8), Inches(2.3), Inches(4.3), Inches(4.7), BG_CARD, RULE)
add_text(sl, Inches(8.3), Inches(2.45), Inches(3.8), Inches(0.25),
         "WHY CONVERSATIONAL?", font_size=9, color=INFO, bold=True, font_name='Consolas')

reasons = [
    ("Mental model match", "Nutritionists think in constraints ('I need X for Y'), not spreadsheet cells. NL matches how they already describe problems."),
    ("Onboarding: minutes, not days", "Legacy tools (Bestmix, Format) require training. FeedIQ works like asking a colleague."),
    ("Knowledge capture", "The AI embeds expert knowledge from 389K analyses. It scales what 8,000 departing employees knew."),
    ("Risk-aware by default", "Every response includes mycotoxin risk — users don't need to remember to check."),
]

ry = Inches(2.8)
for title, desc in reasons:
    add_text(sl, Inches(8.3), ry, Inches(3.8), Inches(0.22),
             title, font_size=9, color=ACCENT, bold=True)
    add_text(sl, Inches(8.3), ry + Inches(0.24), Inches(3.8), Inches(0.8),
             desc, font_size=8, color=DIM)
    ry += Inches(1.0)

add_attribution(sl, "UX concept by Puneet Arora  |  Interaction design: conversational AI + constraint optimization")
add_page_num(sl, 15)


# ══════════════════════════════════════════════════════════════
# SLIDE 16: FEATURE — TRADE-OFF ENGINE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEVELOP", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  FEATURE: TRADE-OFF TRANSPARENCY")
add_title_text(sl, Inches(0.8), Inches(1.0), "AI explains the why,\nnot just the what")

add_body_text(sl, Inches(0.8), Inches(1.9),
              "AI trust comes from transparency. FeedIQ shows users exactly what changes when they\n"
              "adjust budget or risk tolerance — and explains the trade-off in plain language.",
              width=Inches(10), font_size=13)

# Trade-off examples
add_rect(sl, Inches(0.8), Inches(2.9), Inches(5.5), Inches(3.8), BG_CARD, RULE)
add_text(sl, Inches(1.1), Inches(3.05), Inches(5), Inches(0.25),
         "TRADE-OFF EXAMPLES", font_size=9, color=FAINT, bold=True, font_name='Consolas')

tradeoffs = [
    ("CORN SOURCE SWAP", "SE Asian corn ($180/ton, 56% FUM risk)  →  Oceania corn ($192/ton, 8% FUM risk)", "+$12/ton cost   |   -48% fumonisin risk   |   NET: saves $40-60/ton in production losses"),
    ("SOY REDUCTION", "Soybean meal (22%)  →  Rice bran + DDGS blend (15% soy, 10% rice, 5% DDGS)", "-$18/ton cost   |   +12% DON from DDGS   |   Binder added to offset"),
    ("BINDER ADDITION", "No mycotoxin binder  →  0.5% bentonite clay blend", "+$4/ton cost   |   -40% DON absorption   |   ROI: 10-15x in reduced mortality"),
]

ty = Inches(3.4)
for title, change, impact in tradeoffs:
    add_text(sl, Inches(1.1), ty, Inches(5), Inches(0.22),
             title, font_size=9, color=ACCENT, bold=True, font_name='Consolas')
    add_text(sl, Inches(1.1), ty + Inches(0.25), Inches(5), Inches(0.45),
             change, font_size=8, color=DIM)
    add_text(sl, Inches(1.1), ty + Inches(0.68), Inches(5), Inches(0.35),
             impact, font_size=7, color=WHITE, bold=True, font_name='Consolas')
    ty += Inches(1.15)

# Design principle
add_rect(sl, Inches(6.5), Inches(2.9), Inches(5.8), Inches(3.8), BG_CARD, WARN)
add_text(sl, Inches(6.8), Inches(3.05), Inches(5.2), Inches(0.25),
         "DESIGN PRINCIPLE: NO BLACK BOXES", font_size=9, color=WARN, bold=True, font_name='Consolas')

add_text(sl, Inches(6.8), Inches(3.5), Inches(5.2), Inches(3),
         "Most AI products fail not because they're wrong,\n"
         "but because users can't verify they're right.\n\n"
         "FeedIQ makes every decision inspectable:\n\n"
         "• Interactive sliders let users explore alternatives\n"
         "• Cost-risk bars show what changes with each decision\n"
         "• AI explains rationale in plain language\n"
         "• Ingredient-level risk breakdown (not just a score)\n\n"
         "This is 'safety by architecture' — the system\n"
         "is transparent because it was designed to be,\n"
         "not because we added an 'explain' button.\n\n"
         "Trust is the product. The formulation is the proof.",
         font_size=11, color=DIM)

add_attribution(sl, "Puneet Arora  |  Inspired by 20 years of enterprise AI trust patterns at Dell, Boeing, HP")
add_page_num(sl, 16)


# ══════════════════════════════════════════════════════════════
# SLIDE 17: TCT — FROM PRODUCT TO PLATFORM
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEVELOP", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  PLATFORM EVOLUTION")
add_title_text(sl, Inches(0.8), Inches(1.0), "From product to platform:\nThe Cargill Terminal (T-C-T)")

add_body_text(sl, Inches(0.8), Inches(1.9),
              "FeedIQ proved the concept — AI can transform Cargill's data into decisions. T-C-T scales that\n"
              "to an enterprise Decision Intelligence platform spanning three business-critical domains.",
              width=Inches(10), font_size=13)

# Before/After comparison
add_rect(sl, Inches(0.8), Inches(2.9), Inches(5.5), Inches(4.0), BG_CARD, RULE)
add_text(sl, Inches(1.1), Inches(3.05), Inches(5), Inches(0.25),
         "FEEDIQ (PRODUCT)", font_size=9, color=ACCENT, bold=True, font_name='Consolas')
feediq_items = [
    "Single domain: Feed Intelligence",
    "Mycotoxin risk heatmap (8 regions × 6 toxins)",
    "AI feed formulator (4 species scenarios)",
    "Cost-risk trade-off sliders",
    "Chat-based interaction model",
    "~1,600 lines of code",
]
iy = Inches(3.4)
for item in feediq_items:
    add_text(sl, Inches(1.1), iy, Inches(5), Inches(0.28),
             f"  {item}", font_size=10, color=DIM)
    iy += Inches(0.32)

add_rect(sl, Inches(6.5), Inches(2.9), Inches(5.8), Inches(4.0), BG_CARD, INFO)
add_text(sl, Inches(6.8), Inches(3.05), Inches(5.2), Inches(0.25),
         "T-C-T (PLATFORM)", font_size=9, color=INFO, bold=True, font_name='Consolas')
tct_items = [
    "3 domains: Executive + Procurement + Supply Chain",
    "Bloomberg Terminal density UX",
    "AI decision engines with confidence scoring",
    "Supplier Intelligence (8 global suppliers, EUDR)",
    "Shipment tracking + delay prediction AI",
    "Command Center unifying all 3 engines",
    "Intent-based navigation (9 decision modes)",
    "Live commodity ticker + market intelligence",
]
iy = Inches(3.4)
for item in tct_items:
    add_text(sl, Inches(6.8), iy, Inches(5.2), Inches(0.28),
             f"  {item}", font_size=10, color=DIM)
    iy += Inches(0.3)

add_attribution(sl, "Product strategy & architecture by Puneet Arora")
add_page_num(sl, 17)


# ══════════════════════════════════════════════════════════════
# SLIDE 18: TCT — ARCHITECTURE & 3 ENGINES
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEVELOP", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  PLATFORM ARCHITECTURE")
add_title_text(sl, Inches(0.8), Inches(1.0), "Three engines, one intelligence shell")

add_body_text(sl, Inches(0.8), Inches(1.9),
              "T-C-T wraps three deep decision engines in a unified shell inspired by Bloomberg Terminal,\n"
              "Palantir Foundry, and Arc Browser — maximum information density with progressive disclosure.",
              width=Inches(10), font_size=13)

# Shell architecture
add_rect(sl, Inches(0.8), Inches(2.9), Inches(11.5), Inches(0.7), BG_CARD2, ACCENT)
add_text(sl, Inches(1.1), Inches(3.0), Inches(10.5), Inches(0.5),
         "SHELL: Topbar (live ticker + search + clock)  →  Sidebar (engine nav + badges)  →  Main (engine views)  →  Intent Tabs (9 modes)",
         font_size=11, color=DIM)

engines = [
    ("EXECUTIVE ENGINE", "Strategic intelligence for C-suite", ACCENT,
     ["Morning Brief (3 daily priorities)", "Business Pulse (6 KPIs + sparklines)",
      "Strategic Risk tracker (severity + trend)", "Commodity Intelligence (prices + forecasts)",
      "AI Recommendations with confidence scores"]),
    ("PROCUREMENT ENGINE", "Supplier & contract intelligence", INFO,
     ["Supplier Health Scorecard (8 suppliers, 0-10)", "EUDR compliance tracking (PASS/FAIL/PENDING)",
      "Spend Analytics by category ($38.2B)", "Contract Copilot (AI renewal recommendations)",
      "Vendor Risk Assessment + approval pipeline"]),
    ("SUPPLY CHAIN ENGINE", "Logistics & shipment intelligence", WARN,
     ["Mission Control (SVG route map)", "Shipment Intel (6 active, status tracking)",
      "Delay Prediction AI (root cause + ETA)", "Port Intelligence (5 ports, congestion index)",
      "Weather Intelligence + route optimization"]),
]

x = Inches(0.8)
for i, (title, subtitle, color, items) in enumerate(engines):
    cx = x + i * Inches(3.9)
    add_rect(sl, cx, Inches(3.9), Inches(3.7), Inches(3.1), BG_CARD, color)
    add_text(sl, cx + Inches(0.25), Inches(4.0), Inches(3.2), Inches(0.25),
             title, font_size=10, color=color, bold=True, font_name='Consolas')
    add_text(sl, cx + Inches(0.25), Inches(4.3), Inches(3.2), Inches(0.3),
             subtitle, font_size=11, color=WHITE, bold=True)
    iy = Inches(4.7)
    for item in items:
        add_text(sl, cx + Inches(0.25), iy, Inches(3.2), Inches(0.35),
                 f"  {item}", font_size=7, color=DIM)
        iy += Inches(0.32)

add_attribution(sl, "Architecture by Puneet Arora  |  Inspired by Bloomberg Terminal + Palantir Foundry")
add_page_num(sl, 18)


# ══════════════════════════════════════════════════════════════
# SLIDE 19: TCT — EXECUTIVE ENGINE DEEP DIVE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEVELOP", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  EXECUTIVE ENGINE")
add_title_text(sl, Inches(0.8), Inches(1.0), "Every morning starts with\na decision brief")

add_body_text(sl, Inches(0.8), Inches(1.9),
              "The Executive Engine synthesizes overnight events into a prioritized brief.\n"
              "Not a dashboard — a decision queue. Each item has urgency, confidence, and a recommended action.",
              width=Inches(10), font_size=13)

# Morning Brief example
add_rect(sl, Inches(0.8), Inches(2.9), Inches(7.3), Inches(3.8), BG_CARD, ACCENT)
add_text(sl, Inches(1.1), Inches(3.05), Inches(6.5), Inches(0.25),
         "MORNING BRIEF — SAMPLE", font_size=9, color=ACCENT, bold=True, font_name='Consolas')

briefs = [
    ("URGENT", "Soybean port strike probability: 68%", "Recommend: Accelerate Brazil corn sourcing to hedge 5K Asian fumonisin risk", DANGER),
    ("HIGH", "Soybean volatility at 18-month high", "AI: Derivative hedging available — Argentina drought + USDA report speculation drives 3-day trend", WARN),
    ("MONITOR", "Bunge-Viterra integration reaching critical phase", "89% integration completed — South America blending capacity shift could benefit Cargill market share", INFO),
]
by = Inches(3.45)
for urg, title, desc, color in briefs:
    add_rect(sl, Inches(1.0), by, Inches(0.85), Inches(0.22), color, color)
    add_text(sl, Inches(1.05), by + Inches(0.01), Inches(0.8), Inches(0.2),
             urg, font_size=7, color=BG_DARK, bold=True, font_name='Consolas', alignment=PP_ALIGN.CENTER)
    add_text(sl, Inches(2.0), by, Inches(5.8), Inches(0.25),
             title, font_size=9, color=WHITE, bold=True)
    add_text(sl, Inches(2.0), by + Inches(0.28), Inches(5.8), Inches(0.75),
             desc, font_size=7, color=DIM)
    by += Inches(1.1)

# KPIs
add_rect(sl, Inches(8.3), Inches(2.9), Inches(4.0), Inches(3.8), BG_CARD, RULE)
add_text(sl, Inches(8.6), Inches(3.05), Inches(3.5), Inches(0.25),
         "BUSINESS PULSE KPIs", font_size=9, color=ACCENT, bold=True, font_name='Consolas')

kpis = [
    ("$62.48", "CORN FUTURES", "▲ 8.4%"),
    ("2,847", "ACTIVE SHIPMENTS", "▲ 12"),
    ("94.2%", "ON-TIME DELIVERY", "▼ 0.3%"),
    ("7.8", "SUPPLIER HEALTH", "— stable"),
    ("14", "PENDING DECISIONS", "▼ 3"),
]
ky = Inches(3.4)
for val, lbl, trend in kpis:
    add_text(sl, Inches(8.6), ky, Inches(1.5), Inches(0.25),
             val, font_size=12, color=WHITE, bold=True, font_name='Consolas')
    add_text(sl, Inches(10.2), ky, Inches(1.5), Inches(0.22),
             lbl, font_size=7, color=FAINT, font_name='Consolas')
    tcolor = ACCENT if "▲" in trend else DANGER if "▼" in trend else DIM
    add_text(sl, Inches(10.2), ky + Inches(0.22), Inches(1.5), Inches(0.18),
             trend, font_size=8, color=tcolor, font_name='Consolas')
    ky += Inches(0.5)

add_attribution(sl, "Executive intelligence design by Puneet Arora")
add_page_num(sl, 19)


# ══════════════════════════════════════════════════════════════
# SLIDE 20: TCT — PROCUREMENT ENGINE DEEP DIVE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEVELOP", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  PROCUREMENT ENGINE")
add_title_text(sl, Inches(0.8), Inches(1.0), "Supplier intelligence meets\nEUDR compliance")

add_body_text(sl, Inches(0.8), Inches(1.9),
              "Cargill's $38B procurement spend demands AI-native supplier intelligence.\n"
              "The Procurement Engine tracks 8 global suppliers with health scoring, EUDR compliance, and AI contract recommendations.",
              width=Inches(10), font_size=13)

# Supplier scorecard
add_rect(sl, Inches(0.8), Inches(2.9), Inches(7.3), Inches(3.8), BG_CARD, INFO)
add_text(sl, Inches(1.1), Inches(3.05), Inches(6.5), Inches(0.25),
         "SUPPLIER HEALTH SCORECARD (8 GLOBAL SUPPLIERS)", font_size=9, color=INFO, bold=True, font_name='Consolas')

suppliers = [
    ("ADM — Argentina", "8.4", "96.2%", "PASS", "LOW"),
    ("Bunge — Brazil", "8.1", "94.8%", "PASS", "LOW"),
    ("Wilmar — Indonesia", "6.2", "91.2%", "PENDING", "HIGH"),
    ("Louis Dreyfus — France", "7.8", "93.6%", "PASS", "MED"),
    ("COFCO — China", "6.8", "91.5%", "PENDING", "MED"),
    ("Olam — Singapore", "7.5", "93.1%", "PASS", "LOW"),
    ("Glencore Agri — CH", "6.4", "89.8%", "PENDING", "MED"),
    ("Viterra — NL", "9.2", "97.1%", "PASS", "LOW"),
]

add_text(sl, Inches(1.1), Inches(3.35), Inches(1.8), Inches(0.2),
         "SUPPLIER", font_size=8, color=FAINT, bold=True, font_name='Consolas')
add_text(sl, Inches(3.0), Inches(3.35), Inches(0.6), Inches(0.2),
         "HEALTH", font_size=8, color=FAINT, bold=True, font_name='Consolas')
add_text(sl, Inches(3.7), Inches(3.35), Inches(0.8), Inches(0.2),
         "ON-TIME", font_size=8, color=FAINT, bold=True, font_name='Consolas')
add_text(sl, Inches(4.6), Inches(3.35), Inches(0.8), Inches(0.2),
         "EUDR", font_size=8, color=FAINT, bold=True, font_name='Consolas')
add_text(sl, Inches(5.5), Inches(3.35), Inches(0.8), Inches(0.2),
         "RISK", font_size=8, color=FAINT, bold=True, font_name='Consolas')

sy = Inches(3.6)
for name, health, ontime, eudr, risk in suppliers:
    add_text(sl, Inches(1.1), sy, Inches(1.8), Inches(0.2),
             name, font_size=8, color=DIM, font_name='Consolas')
    hcolor = ACCENT if float(health) >= 7.5 else WARN if float(health) >= 6.5 else DANGER
    add_text(sl, Inches(3.0), sy, Inches(0.6), Inches(0.2),
             health, font_size=9, color=hcolor, bold=True, font_name='Consolas')
    add_text(sl, Inches(3.7), sy, Inches(0.8), Inches(0.2),
             ontime, font_size=8, color=DIM, font_name='Consolas')
    ecolor = ACCENT if eudr == "PASS" else WARN
    add_text(sl, Inches(4.6), sy, Inches(0.8), Inches(0.2),
             eudr, font_size=8, color=ecolor, bold=True, font_name='Consolas')
    rcolor = ACCENT if risk == "LOW" else WARN if risk == "MED" else DANGER
    add_text(sl, Inches(5.5), sy, Inches(0.8), Inches(0.2),
             risk, font_size=8, color=rcolor, bold=True, font_name='Consolas')
    sy += Inches(0.28)

# Contract Copilot
add_rect(sl, Inches(8.3), Inches(2.9), Inches(4.0), Inches(3.8), BG_CARD, RULE)
add_text(sl, Inches(8.6), Inches(3.05), Inches(3.5), Inches(0.25),
         "CONTRACT COPILOT AI", font_size=9, color=ACCENT, bold=True, font_name='Consolas')

add_text(sl, Inches(8.6), Inches(3.5), Inches(3.5), Inches(0.25),
         "Wilmar renewal due in 14 days", font_size=10, color=WHITE, bold=True)
add_text(sl, Inches(8.6), Inches(3.8), Inches(3.5), Inches(1.0),
         "AI suggests: renegotiate with EUDR\ncompliance clause. Add 60-day GPS\npolygon mapping requirement.\nCompeting palm suppliers offer\n3% lower FOB.",
         font_size=8, color=DIM)

add_text(sl, Inches(8.6), Inches(5.0), Inches(3.5), Inches(0.25),
         "ADM contract optimization", font_size=9, color=WHITE, bold=True)
add_text(sl, Inches(8.6), Inches(5.28), Inches(3.5), Inches(0.8),
         "Volume commitment increase of 8%\nunlocks tier-2 pricing. Estimated\nsaving: $4.2M annually.\nAI confidence: 89%",
         font_size=7, color=DIM)

add_rect(sl, Inches(8.6), Inches(6.25), Inches(1.0), Inches(0.25), ACCENT, ACCENT)
add_text(sl, Inches(8.65), Inches(6.27), Inches(0.9), Inches(0.2),
         "DRAFT", font_size=7, color=BG_DARK, bold=True, font_name='Consolas', alignment=PP_ALIGN.CENTER)
add_rect(sl, Inches(9.8), Inches(6.25), Inches(1.0), Inches(0.25), INFO, INFO)
add_text(sl, Inches(9.85), Inches(6.27), Inches(0.9), Inches(0.2),
         "SIMULATE", font_size=7, color=BG_DARK, bold=True, font_name='Consolas', alignment=PP_ALIGN.CENTER)

add_attribution(sl, "Procurement intelligence design by Puneet Arora")
add_page_num(sl, 20)


# ══════════════════════════════════════════════════════════════
# SLIDE 21: TCT — SUPPLY CHAIN ENGINE DEEP DIVE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEVELOP", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  SUPPLY CHAIN ENGINE")
add_title_text(sl, Inches(0.8), Inches(1.0), "Predict disruptions before\nthey hit the bottom line")

add_body_text(sl, Inches(0.8), Inches(1.9),
              "2,847 active shipments, 94.2% on-time delivery, 5 critical ports — the Supply Chain Engine\n"
              "gives operators predictive intelligence to reroute, pre-position, and mitigate before delays compound.",
              width=Inches(10), font_size=13)

# Shipment table
add_rect(sl, Inches(0.8), Inches(2.9), Inches(7.3), Inches(2.0), BG_CARD, WARN)
add_text(sl, Inches(1.1), Inches(3.05), Inches(6.5), Inches(0.25),
         "ACTIVE SHIPMENTS (SAMPLE)", font_size=9, color=WARN, bold=True, font_name='Consolas')

add_text(sl, Inches(1.1), Inches(3.35), Inches(1.0), Inches(0.2),
         "ID", font_size=8, color=FAINT, bold=True, font_name='Consolas')
add_text(sl, Inches(2.0), Inches(3.35), Inches(1.5), Inches(0.2),
         "ROUTE", font_size=8, color=FAINT, bold=True, font_name='Consolas')
add_text(sl, Inches(3.7), Inches(3.35), Inches(1.0), Inches(0.2),
         "CARGO", font_size=8, color=FAINT, bold=True, font_name='Consolas')
add_text(sl, Inches(4.8), Inches(3.35), Inches(0.8), Inches(0.2),
         "STATUS", font_size=8, color=FAINT, bold=True, font_name='Consolas')
add_text(sl, Inches(5.8), Inches(3.35), Inches(1.2), Inches(0.2),
         "ETA", font_size=8, color=FAINT, bold=True, font_name='Consolas')
add_text(sl, Inches(6.8), Inches(3.35), Inches(1.0), Inches(0.2),
         "DELAY", font_size=8, color=FAINT, bold=True, font_name='Consolas')

shipments = [
    ("CGL-8891", "Santos → Rotterdam", "Soybean", "IN TRANSIT", "Jun 28", "0d"),
    ("CGL-8823", "Odessa → Jakarta", "Wheat", "DELAYED", "Jul 03", "+3d"),
    ("CGL-8756", "NOLA → Shanghai", "Corn", "IN TRANSIT", "Jul 12", "0d"),
    ("CGL-8712", "Rosario → Mumbai", "Meal", "DELAYED", "Jul 05", "+1d"),
]

ty = Inches(3.6)
for sid, route, cargo, status, eta, delay in shipments:
    add_text(sl, Inches(1.1), ty, Inches(1.0), Inches(0.2),
             sid, font_size=8, color=DIM, font_name='Consolas')
    add_text(sl, Inches(2.0), ty, Inches(1.5), Inches(0.2),
             route, font_size=8, color=DIM, font_name='Consolas')
    add_text(sl, Inches(3.7), ty, Inches(1.0), Inches(0.2),
             cargo, font_size=8, color=DIM, font_name='Consolas')
    scolor = WARN if status == "DELAYED" else ACCENT
    add_text(sl, Inches(4.8), ty, Inches(0.8), Inches(0.2),
             status, font_size=7, color=scolor, bold=True, font_name='Consolas')
    add_text(sl, Inches(5.8), ty, Inches(1.0), Inches(0.2),
             eta, font_size=8, color=DIM, font_name='Consolas')
    dcolor = DANGER if "+" in delay else ACCENT
    add_text(sl, Inches(6.8), ty, Inches(1.0), Inches(0.2),
             delay, font_size=8, color=dcolor, bold=True, font_name='Consolas')
    ty += Inches(0.26)

# Delay prediction AI
add_rect(sl, Inches(0.8), Inches(5.1), Inches(3.5), Inches(1.7), BG_CARD, DANGER)
add_text(sl, Inches(1.1), Inches(5.25), Inches(3), Inches(0.25),
         "DELAY PREDICTION AI", font_size=9, color=DANGER, bold=True, font_name='Consolas')
add_text(sl, Inches(1.1), Inches(5.55), Inches(3), Inches(0.25),
         "CGL-8823: 3-day delay predicted", font_size=9, color=WHITE, bold=True)
add_text(sl, Inches(1.1), Inches(5.82), Inches(3), Inches(0.85),
         "Root cause: Odessa port congestion\n(32 vessels, 4.2d avg dwell).\nAI confidence: 87%\nRecommend: Reroute via Constanta.",
         font_size=7, color=DIM)

# Port intelligence
add_rect(sl, Inches(4.5), Inches(5.1), Inches(3.6), Inches(1.7), BG_CARD, RULE)
add_text(sl, Inches(4.8), Inches(5.25), Inches(3), Inches(0.25),
         "PORT INTELLIGENCE (TOP 5)", font_size=9, color=INFO, bold=True, font_name='Consolas')
ports = [
    ("Santos", "18 vessels", "LOW"),
    ("Rotterdam", "24 vessels", "MED"),
    ("Odessa", "32 vessels", "HIGH"),
    ("NOLA", "15 vessels", "LOW"),
    ("Shanghai", "28 vessels", "MED"),
]
py = Inches(5.55)
for port, vessels, cong in ports:
    add_text(sl, Inches(4.8), py, Inches(1.3), Inches(0.22),
             port, font_size=8, color=DIM, font_name='Consolas')
    add_text(sl, Inches(6.0), py, Inches(1.0), Inches(0.22),
             vessels, font_size=8, color=DIM, font_name='Consolas')
    ccolor = ACCENT if cong == "LOW" else WARN if cong == "MED" else DANGER
    add_text(sl, Inches(7.0), py, Inches(0.8), Inches(0.22),
             cong, font_size=7, color=ccolor, bold=True, font_name='Consolas')
    py += Inches(0.24)

# KPI cards
add_rect(sl, Inches(8.3), Inches(2.9), Inches(4.0), Inches(3.9), BG_CARD, RULE)
add_text(sl, Inches(8.6), Inches(3.05), Inches(3.5), Inches(0.25),
         "SUPPLY CHAIN KPIs", font_size=9, color=ACCENT, bold=True, font_name='Consolas')

sc_kpis = [
    ("2,847", "Active shipments"),
    ("94.2%", "On-time delivery"),
    ("38", "Global routes active"),
    ("2.46", "Avg transit ratio"),
    ("3", "Active weather alerts"),
]
ky = Inches(3.45)
for val, lbl in sc_kpis:
    add_text(sl, Inches(8.6), ky, Inches(2), Inches(0.28),
             val, font_size=16, color=WHITE, bold=True, font_name='Consolas')
    add_text(sl, Inches(8.6), ky + Inches(0.3), Inches(3), Inches(0.2),
             lbl, font_size=8, color=FAINT, font_name='Consolas')
    ky += Inches(0.58)

add_attribution(sl, "Supply chain intelligence design by Puneet Arora")
add_page_num(sl, 21)


# ══════════════════════════════════════════════════════════════
# SLIDE 22: TCT — INTENT-BASED UX & AI CONFIDENCE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEVELOP", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  UX INNOVATION: INTENT-BASED NAVIGATION")
add_title_text(sl, Inches(0.8), Inches(1.0), "9 intent modes map the\nentire decision lifecycle")

add_body_text(sl, Inches(0.8), Inches(1.9),
              "Traditional dashboards organize by data type. T-C-T organizes by user intent —\n"
              "what you're trying to DO, not what you're looking AT. This is the core UX innovation.",
              width=Inches(10), font_size=13)

# Intent lifecycle
intents = [
    ("OBSERVE", "Monitor data streams,\nspot anomalies", ACCENT),
    ("UNDERSTAND", "Contextualize signals\nwith historical data", ACCENT),
    ("INVESTIGATE", "Drill into root causes\nand dependencies", INFO),
    ("PREDICT", "AI-generated forecasts\nwith confidence bands", INFO),
    ("RECOMMEND", "Ranked options with\ntrade-off analysis", WARN),
    ("SIMULATE", "What-if scenarios with\nimpact projections", WARN),
    ("APPROVE", "Decision queue with\naudit trail + sign-off", DANGER),
    ("EXECUTE", "One-click action with\nrollback capability", DANGER),
    ("LEARN", "Outcome tracking feeds\nback into AI models", FAINT),
]

ix = Inches(0.6)
iw = Inches(1.33)
for i, (name, desc, color) in enumerate(intents):
    cx = ix + i * (iw + Inches(0.05))
    add_rect(sl, cx, Inches(2.9), iw, Inches(1.4), BG_CARD, color)
    add_text(sl, cx + Inches(0.08), Inches(2.98), iw - Inches(0.16), Inches(0.2),
             name, font_size=7, color=color, bold=True, font_name='Consolas', alignment=PP_ALIGN.CENTER)
    add_text(sl, cx + Inches(0.08), Inches(3.22), iw - Inches(0.16), Inches(0.9),
             desc, font_size=6, color=DIM, alignment=PP_ALIGN.CENTER)
    if i < 8:
        add_text(sl, cx + iw, Inches(3.3), Inches(0.1), Inches(0.2),
                 "→", font_size=8, color=FAINT, alignment=PP_ALIGN.CENTER)

# AI Confidence
add_rect(sl, Inches(0.8), Inches(4.7), Inches(5.5), Inches(2.2), BG_CARD, ACCENT)
add_text(sl, Inches(1.1), Inches(4.82), Inches(5), Inches(0.25),
         "AI RECOMMENDATION CARDS", font_size=9, color=ACCENT, bold=True, font_name='Consolas')
add_text(sl, Inches(1.1), Inches(5.1), Inches(5), Inches(1.7),
         "Every AI recommendation in T-C-T includes:\n\n"
         "• Confidence score (0-100%) with visual meter\n"
         "• Urgency badge (CRITICAL / HIGH / MONITOR)\n"
         "• Estimated impact (cost savings, risk reduction)\n"
         "• Action buttons (APPROVE / SIMULATE / INVESTIGATE)\n"
         "• Decision audit trail for compliance",
         font_size=8, color=DIM)

# Comparison with traditional
add_rect(sl, Inches(6.5), Inches(4.7), Inches(5.8), Inches(2.2), BG_CARD, RULE)
add_text(sl, Inches(6.8), Inches(4.82), Inches(5.2), Inches(0.25),
         "TRADITIONAL vs. INTENT-BASED UX", font_size=9, color=INFO, bold=True, font_name='Consolas')

comparisons = [
    ("Dashboard", "Shows data, user decides what to do", DIM),
    ("Chatbot", "Answers questions, no workflow context", DIM),
    ("Alert System", "Flags issues, no resolution path", DIM),
    ("T-C-T", "Guides the entire decision lifecycle\nfrom observation through execution", ACCENT),
]
cy = Inches(5.15)
for label, desc, color in comparisons:
    is_tct = label == "T-C-T"
    if is_tct:
        add_rect(sl, Inches(6.7), cy - Inches(0.05), Inches(5.4), Inches(0.7), BG_CARD2, ACCENT)
    add_text(sl, Inches(6.8), cy, Inches(1.2), Inches(0.22),
             label, font_size=9, color=ACCENT if is_tct else FAINT, bold=True, font_name='Consolas')
    add_text(sl, Inches(8.0), cy, Inches(4), Inches(0.5),
             desc, font_size=8, color=WHITE if is_tct else DIM)
    cy += Inches(0.45) if not is_tct else Inches(0.7)

add_attribution(sl, "Intent-based UX framework by Puneet Arora  |  Informed by Apple Intelligence + Arc Browser")
add_page_num(sl, 22)


# ══════════════════════════════════════════════════════════════
# SLIDE 23: DESIGN PRINCIPLES
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEVELOP", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  UX DESIGN PRINCIPLES")
add_title_text(sl, Inches(0.8), Inches(1.0), "Three principles that shaped\nevery FeedIQ decision")

principles = [
    ("01", "CONVERSATIONAL\nOVER FORMS", ACCENT,
     "Feed nutritionists think in constraints,\nnot spreadsheet cells.",
     "A natural language interface matches\ntheir mental model: 'I need X for Y\nconditions' — not 'fill row 47, col C.'\n\nReduces onboarding from days to minutes.\nCaptures expert knowledge in every query."),
    ("02", "RISK-AWARE\nBY DEFAULT", INFO,
     "Every formulation automatically\nchecks mycotoxin risk.",
     "This isn't a separate workflow — it's\nembedded. Users don't need to remember\nto check; the system checks for them.\n\n'Safety by architecture' — the most\nreliable safety net is the one that\nrequires zero user action."),
    ("03", "TRADE-OFF\nTRANSPARENCY", WARN,
     "AI systems fail when they're\nblack boxes.",
     "FeedIQ shows the WHY — cost-risk\nsliders let users explore alternatives,\nsee what changes, and build trust.\n\nInspired by Cargill's own 'explain\nthe trade-off' philosophy. Trust is\nearned through transparency, not claims."),
]

x = Inches(0.8)
for i, (num, title, color, hook, detail) in enumerate(principles):
    cx = x + i * Inches(3.9)
    add_rect(sl, cx, Inches(2.4), Inches(3.7), Inches(4.7), BG_CARD, color)
    add_text(sl, cx + Inches(0.25), Inches(2.55), Inches(0.5), Inches(0.35),
             num, font_size=24, color=color, bold=True, font_name='Consolas')
    add_text(sl, cx + Inches(0.25), Inches(2.95), Inches(3.2), Inches(0.6),
             title, font_size=14, color=WHITE, bold=True)
    add_text(sl, cx + Inches(0.25), Inches(3.6), Inches(3.2), Inches(0.65),
             hook, font_size=9, color=color)
    add_divider(sl, cx + Inches(0.25), Inches(4.3), Inches(3.2))
    add_text(sl, cx + Inches(0.25), Inches(4.45), Inches(3.2), Inches(2.5),
             detail, font_size=8, color=DIM)

add_attribution(sl, "Design principles by Puneet Arora  |  Informed by Dell GenAI VA (4.4% → 80% engagement)")
add_page_num(sl, 23)


# ══════════════════════════════════════════════════════════════
# SLIDE 24: INFORMATION ARCHITECTURE
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DEVELOP", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  INFORMATION ARCHITECTURE")
add_title_text(sl, Inches(0.8), Inches(1.0), "User flow: question to optimized feed")

# Flow steps
steps = [
    ("INPUT", "User describes need in\nnatural language\n\n'Broiler feed, tropical,\n$280/ton, low soy'", ACCENT),
    ("PARSE", "LLM extracts:\n• Species + stage\n• Climate zone\n• Budget constraint\n• Ingredient preferences", INFO),
    ("OPTIMIZE", "Constraint solver:\n• 15+ nutritional params\n• Cost minimization\n• Local ingredient DB\n• Availability check", WARN),
    ("RISK CHECK", "Mycotoxin overlay:\n• Region × toxin scoring\n• Multi-toxin compound risk\n• Threshold comparison\n• Auto-swap if HIGH", DANGER),
    ("PRESENT", "Formulation output:\n• Ingredient breakdown\n• Donut + table viz\n• Nutrition profile\n• Risk gauge per item", ACCENT),
    ("ITERATE", "Trade-off sliders:\n• Budget ↔ Risk\n• AI recommendations\n• Source alternatives\n• Export spec sheet", INFO),
]

sx = Inches(0.8)
sy = Inches(2.2)
sw = Inches(1.9)
for i, (title, desc, color) in enumerate(steps):
    cx = sx + i * (sw + Inches(0.08))
    add_rect(sl, cx, sy, sw, Inches(3.5), BG_CARD, color)
    add_text(sl, cx + Inches(0.15), sy + Inches(0.1), sw - Inches(0.3), Inches(0.2),
             f"STEP {i+1}", font_size=8, color=FAINT, bold=True, font_name='Consolas')
    add_text(sl, cx + Inches(0.15), sy + Inches(0.35), sw - Inches(0.3), Inches(0.3),
             title, font_size=12, color=color, bold=True, font_name='Consolas')
    add_text(sl, cx + Inches(0.15), sy + Inches(0.75), sw - Inches(0.3), Inches(2.5),
             desc, font_size=8, color=DIM)

# Arrow indicators between steps
for i in range(5):
    ax = sx + (i + 1) * (sw + Inches(0.08)) - Inches(0.06)
    add_text(sl, ax, Inches(3.5), Inches(0.12), Inches(0.3),
             "→", font_size=14, color=FAINT, alignment=PP_ALIGN.CENTER)

# Technical note
add_rect(sl, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.9), BG_CARD, RULE)
add_text(sl, Inches(1.1), Inches(6.0), Inches(10.5), Inches(0.7),
         "TECH STACK: Static HTML/CSS/JS prototype (no build step) → Production: Next.js + Claude API + Linear programming solver + "
         "Cargill mycotoxin DB integration + Supabase for supplier/formulation data. Designed for progressive enhancement — "
         "prototype demonstrates full UX with mock data, production adds real-time data feeds.",
         font_size=10, color=DIM)

add_attribution(sl, "Architecture by Puneet Arora")
add_page_num(sl, 24)


# ══════════════════════════════════════════════════════════════
# SLIDE 25: PROTOTYPE DEMO (DELIVER)
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DELIVER", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  WORKING PROTOTYPES")
add_title_text(sl, Inches(0.8), Inches(1.0), "Two live prototypes: try them now")

# FeedIQ URL
add_rect(sl, Inches(0.8), Inches(2.1), Inches(5.5), Inches(0.9), BG_CARD, ACCENT)
add_text(sl, Inches(1.1), Inches(2.2), Inches(5), Inches(0.25),
         "FEEDIQ — LIVE", font_size=9, color=ACCENT, bold=True, font_name='Consolas')
add_text(sl, Inches(1.1), Inches(2.5), Inches(5), Inches(0.3),
         "puneetindeelhi940.github.io/portfolio/feediq.html", font_size=11, color=WHITE, bold=True, font_name='Consolas')

# TCT URL
add_rect(sl, Inches(6.5), Inches(2.1), Inches(5.8), Inches(0.9), BG_CARD, INFO)
add_text(sl, Inches(6.8), Inches(2.2), Inches(5.2), Inches(0.25),
         "T-C-T — LIVE", font_size=9, color=INFO, bold=True, font_name='Consolas')
add_text(sl, Inches(6.8), Inches(2.5), Inches(5.2), Inches(0.3),
         "puneetindeelhi940.github.io/portfolio/tct.html", font_size=11, color=WHITE, bold=True, font_name='Consolas')

# FeedIQ features
add_rect(sl, Inches(0.8), Inches(3.3), Inches(5.5), Inches(3.4), BG_CARD, RULE)
add_text(sl, Inches(1.1), Inches(3.45), Inches(5), Inches(0.25),
         "FEEDIQ PROTOTYPE", font_size=9, color=ACCENT, bold=True, font_name='Consolas')
feediq_demo = [
    "Global contamination heatmap (8 regions × 6 toxins)",
    "Conversational AI formulator (4 species scenarios)",
    "Ingredient donut chart + cost breakdown table",
    "Nutrition profile with target range bars",
    "Per-ingredient mycotoxin risk assessment",
    "Cost-risk trade-off sliders + AI recommendations",
]
iy = Inches(3.8)
for item in feediq_demo:
    add_text(sl, Inches(1.1), iy, Inches(5), Inches(0.25),
             f"✓  {item}", font_size=10, color=DIM)
    iy += Inches(0.28)

# TCT features
add_rect(sl, Inches(6.5), Inches(3.3), Inches(5.8), Inches(3.5), BG_CARD, RULE)
add_text(sl, Inches(6.8), Inches(3.45), Inches(5.2), Inches(0.25),
         "T-C-T PROTOTYPE", font_size=9, color=INFO, bold=True, font_name='Consolas')
tct_demo = [
    "Command Center with cross-engine KPI dashboard",
    "Executive Engine: Morning Brief + Business Pulse",
    "Procurement Engine: 8-supplier scorecard + EUDR",
    "Supply Chain Engine: Route map + delay prediction",
    "AI recommendation cards with confidence scoring",
    "Intent-based navigation (9 decision modes)",
    "Live commodity ticker (10 instruments)",
    "Bloomberg Terminal-density data layout",
]
iy = Inches(3.8)
for item in tct_demo:
    add_text(sl, Inches(6.8), iy, Inches(5.2), Inches(0.25),
             f"✓  {item}", font_size=10, color=DIM)
    iy += Inches(0.28)

add_attribution(sl, "Both prototypes designed & built by Puneet Arora  |  Vibe-coded in 3 days")
add_page_num(sl, 25)


# ══════════════════════════════════════════════════════════════
# SLIDE 26: METRICS FRAMEWORK
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DELIVER", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  METRICS FRAMEWORK")
add_title_text(sl, Inches(0.8), Inches(1.0), "How we measure success")

# Metric categories
add_rect(sl, Inches(0.8), Inches(2.2), Inches(11.5), Inches(4.5), BG_CARD, RULE)

categories = [
    ("ADOPTION", ACCENT, [
        ("Time to first formulation", "Target: < 2 min", "Benchmark: Bestmix avg setup = 4-6 hours"),
        ("Weekly active users", "Target: 60% of nutritionists", "Benchmark: 30% for legacy tools"),
        ("Prompt-to-formulation rate", "Target: 85%+ completion", "Benchmark: 55% form completion"),
    ]),
    ("QUALITY", INFO, [
        ("Formulation accuracy", "Target: 95% within NRC ranges", "Benchmark: expert nutritionist = 92%"),
        ("Risk detection rate", "Target: 100% HIGH risk flagged", "Benchmark: manual checking = ~60%"),
        ("User trust score", "Target: 8+/10 on trust survey", "Benchmark: avg AI tool = 5.2/10"),
    ]),
    ("BUSINESS IMPACT", WARN, [
        ("Feed cost reduction", "Target: 8-12%", "Benchmark: Adisseo = 12% for pig farms"),
        ("Mycotoxin-related losses", "Target: 30% reduction", "Benchmark: binder alone = 15-20%"),
        ("Knowledge retention", "Target: capture 90% of SOPs", "Benchmark: 0% digital capture today"),
    ]),
]

x = Inches(1.0)
for i, (cat, color, metrics) in enumerate(categories):
    cx = x + i * Inches(3.8)
    add_text(sl, cx, Inches(2.35), Inches(3.4), Inches(0.25),
             cat, font_size=10, color=color, bold=True, font_name='Consolas')
    my = Inches(2.7)
    for metric, target, bench in metrics:
        add_text(sl, cx, my, Inches(3.4), Inches(0.25),
                 metric, font_size=9, color=WHITE, bold=True)
        add_text(sl, cx, my + Inches(0.28), Inches(3.4), Inches(0.25),
                 target, font_size=8, color=color, font_name='Consolas')
        add_text(sl, cx, my + Inches(0.55), Inches(3.4), Inches(0.3),
                 bench, font_size=7, color=FAINT, font_name='Consolas')
        my += Inches(0.95)

add_text(sl, Inches(0.8), Inches(6.85), Inches(11.5), Inches(0.3),
         "Benchmarks sourced from: Adisseo AI platform (2026), Cargill VIV Asia 2025 reports, industry averages",
         font_size=8, color=FAINT, font_name='Consolas')
add_attribution(sl, "Metrics framework by Puneet Arora")
add_page_num(sl, 26)


# ══════════════════════════════════════════════════════════════
# SLIDE 27: PRODUCT ROADMAP
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DELIVER", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  PRODUCT ROADMAP")
add_title_text(sl, Inches(0.8), Inches(1.0), "Three phases to production")

phases = [
    ("PHASE 1", "0 – 3 MONTHS", "Prove the concept", ACCENT,
     ["FeedIQ prototype (DONE)", "T-C-T 3-engine prototype (DONE)", "User testing with 5-10 stakeholders",
      "Validate formulation accuracy vs. expert", "Cargill internal stakeholder review", "Refine AI prompts and engine logic"]),
    ("PHASE 2", "3 – 9 MONTHS", "Connect real data", INFO,
     ["Integrate mycotoxin DB + procurement APIs", "Connect to live commodity & shipment feeds",
      "Claude/GPT API for live AI recommendations", "EUDR compliance data feed integration",
      "Pilot with 2-3 Cargill business units", "A/B test intent-based vs. dashboard UX"]),
    ("PHASE 3", "9 – 18 MONTHS", "Scale & monetize", WARN,
     ["Production deployment across Cargill", "T-C-T as Cargill's internal decision OS",
      "Mobile companion app for field/exec use", "Cross-engine AI decision correlation",
      "Integration with CMAX, CNS, Galleon", "White-label for Cargill's B2B customers"]),
]

x = Inches(0.8)
for i, (phase, time, subtitle, color, items) in enumerate(phases):
    cx = x + i * Inches(3.9)
    add_rect(sl, cx, Inches(2.2), Inches(3.7), Inches(4.9), BG_CARD, color)
    add_text(sl, cx + Inches(0.25), Inches(2.3), Inches(2), Inches(0.25),
             phase, font_size=10, color=color, bold=True, font_name='Consolas')
    add_text(sl, cx + Inches(0.25), Inches(2.6), Inches(3.2), Inches(0.25),
             time, font_size=9, color=FAINT, font_name='Consolas')
    add_text(sl, cx + Inches(0.25), Inches(2.9), Inches(3.2), Inches(0.35),
             subtitle, font_size=16, color=WHITE, bold=True)
    iy = Inches(3.4)
    for item in items:
        prefix = "✓" if "DONE" in item else "○"
        c = ACCENT if "DONE" in item else DIM
        add_text(sl, cx + Inches(0.25), iy, Inches(3.2), Inches(0.38),
                 f" {prefix}  {item}", font_size=8, color=c)
        iy += Inches(0.38)

add_attribution(sl, "Roadmap by Puneet Arora")
add_page_num(sl, 27)


# ══════════════════════════════════════════════════════════════
# SLIDE 28: BENCHMARKS & VALIDATION
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_phase_label(sl, Inches(0.8), Inches(0.5), "DELIVER", 2)
add_label(sl, Inches(2.6), Inches(0.53), "—  INDUSTRY BENCHMARKS")
add_title_text(sl, Inches(0.8), Inches(1.0), "Similar AI tools validate\nthe FeedIQ approach")

benchmarks = [
    ("Adisseo AI Feed Platform", "March 2026", "12% cost reduction", "AI-driven formulation for Chinese pig farms. Validates that AI feed optimization delivers measurable savings. FeedIQ adds: global mycotoxin risk overlay.", ACCENT),
    ("Cargill Port Optimizer (CMAX)", "2024-2025", "30x ROI", "Cargill's own AI for shipping logistics. Proves Cargill can operationalize AI. FeedIQ extends this to the feed formulation domain.", INFO),
    ("Cargill Manufacturing Analytics", "2024-2025", "$15M+ in benefits", "AI-driven process optimization across manufacturing. Demonstrates data-to-action value. FeedIQ applies the same principle to feed data.", WARN),
    ("Dell GenAI Virtual Assistant", "2024-2026", "4.4% → 80% engagement", "Puneet Arora's prior work at Dell — transformed a low-engagement chatbot into an 80% engagement, +14% conversion AI tool. Same design principles applied to FeedIQ.", DANGER),
]

by = Inches(2.2)
for title, date, metric, desc, color in benchmarks:
    add_rect(sl, Inches(0.8), by, Inches(11.5), Inches(0.95), BG_CARD, RULE)
    add_text(sl, Inches(1.1), by + Inches(0.06), Inches(3), Inches(0.22),
             title, font_size=10, color=WHITE, bold=True)
    add_text(sl, Inches(1.1), by + Inches(0.32), Inches(3), Inches(0.2),
             date, font_size=8, color=FAINT, font_name='Consolas')
    add_text(sl, Inches(4.2), by + Inches(0.08), Inches(1.8), Inches(0.35),
             metric, font_size=13, color=color, bold=True, font_name='Consolas')
    add_text(sl, Inches(6.2), by + Inches(0.06), Inches(5.8), Inches(0.85),
             desc, font_size=8, color=DIM)
    by += Inches(1.05)

add_rect(sl, Inches(0.8), Inches(6.45), Inches(11.5), Inches(0.55), BG_CARD2, ACCENT)
add_text(sl, Inches(1.1), Inches(6.5), Inches(10.5), Inches(0.45),
         "PATTERN: AI tools that connect domain data to user decisions consistently deliver 10-30x ROI. "
         "FeedIQ applies this proven pattern to a $1.5B market with zero existing interactive tools.",
         font_size=10, color=DIM)

add_attribution(sl, "Research & analysis by Puneet Arora")
add_page_num(sl, 28)


# ══════════════════════════════════════════════════════════════
# SLIDE 29: SOLUTION NOVELTY
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_label(sl, Inches(0.8), Inches(0.6), "—  SOLUTION NOVELTY")
add_title_text(sl, Inches(0.8), Inches(1.0), "What makes this different")

add_rect(sl, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.0), BG_CARD, RULE)
add_text(sl, Inches(1.1), Inches(2.1), Inches(5), Inches(0.25),
         "FEEDIQ FIRSTS", font_size=9, color=ACCENT, bold=True, font_name='Consolas')

feediq_novel = [
    ("Risk + Formulation unified", "First tool combining mycotoxin risk with feed optimization in one interface.", ACCENT),
    ("Natural language formulation", "Conversational AI replaces form-based UIs. Onboarding: days → minutes.", INFO),
    ("Ingredient-level risk scoring", "Per-ingredient mycotoxin scoring by sourcing region, not just aggregates.", WARN),
    ("Cost-risk trade-off viz", "'Spend $12 more/ton → 48% less fumonisin risk.' Invisible trade-offs made visible.", DANGER),
]

ny = Inches(2.45)
for title, desc, color in feediq_novel:
    add_text(sl, Inches(1.1), ny, Inches(5), Inches(0.25),
             title, font_size=10, color=color, bold=True)
    add_text(sl, Inches(1.1), ny + Inches(0.28), Inches(5), Inches(0.6),
             desc, font_size=9, color=DIM)
    ny += Inches(0.85)

add_rect(sl, Inches(6.5), Inches(2.0), Inches(5.8), Inches(4.0), BG_CARD, RULE)
add_text(sl, Inches(6.8), Inches(2.1), Inches(5.2), Inches(0.25),
         "T-C-T FIRSTS", font_size=9, color=INFO, bold=True, font_name='Consolas')

tct_novel = [
    ("Intent-based enterprise UX", "Navigation by user intent (Observe → Execute), not data type. New UX paradigm.", INFO),
    ("3-engine decision platform", "Executive + Procurement + Supply Chain in one shell. Zero existing parallels at Cargill.", ACCENT),
    ("AI confidence scoring", "Every AI recommendation shows confidence %, urgency, and impact. Explainable AI by design.", WARN),
    ("EUDR compliance intelligence", "Supplier-level deforestation regulation tracking with AI renegotiation recommendations.", DANGER),
]

ny = Inches(2.45)
for title, desc, color in tct_novel:
    add_text(sl, Inches(6.8), ny, Inches(5.2), Inches(0.25),
             title, font_size=10, color=color, bold=True)
    add_text(sl, Inches(6.8), ny + Inches(0.28), Inches(5.2), Inches(0.6),
             desc, font_size=9, color=DIM)
    ny += Inches(0.85)

add_attribution(sl, "Product strategy by Puneet Arora")
add_page_num(sl, 29)


# ══════════════════════════════════════════════════════════════
# SLIDE 30: CLOSING / ABOUT
# ══════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
set_bg(sl)
add_rect(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT)

add_text(sl, Inches(0.8), Inches(1.5), Inches(11), Inches(1.5),
         "From data to decision.\nFrom risk to action.\nFrom complexity to clarity.",
         font_size=36, color=WHITE, bold=True)

add_text(sl, Inches(0.8), Inches(3.5), Inches(8), Inches(0.5),
         "FeedIQ — a product concept demonstrating how AI and thoughtful UX\n"
         "can transform Cargill's most valuable data asset into a competitive advantage.",
         font_size=15, color=DIM)

add_divider(sl, Inches(0.8), Inches(4.5), Inches(11.5))

add_text(sl, Inches(0.8), Inches(4.8), Inches(5), Inches(0.3),
         "PUNEET ARORA", font_size=18, color=WHITE, bold=True, font_name='Consolas')
add_text(sl, Inches(0.8), Inches(5.2), Inches(6), Inches(1.5),
         "Principal Product Designer & AI Design Leader\n"
         "20 years  |  Dell · Boeing · Software AG · HP R&D\n\n"
         "Led Dell's first GenAI Virtual Assistant (4.4% → 80% engagement)\n"
         "33 designers across 5 time zones  |  8 patents granted\n"
         "President & Game Changer Awards at Dell — 2025",
         font_size=12, color=DIM)

add_rect(sl, Inches(7.5), Inches(4.8), Inches(5), Inches(2.2), BG_CARD, RULE)
add_text(sl, Inches(7.8), Inches(4.95), Inches(4.5), Inches(0.25),
         "LINKS & CONTACT", font_size=9, color=ACCENT, bold=True, font_name='Consolas')
links = [
    ("FeedIQ Prototype", "puneetindeelhi940.github.io/portfolio/feediq.html"),
    ("T-C-T Prototype", "puneetindeelhi940.github.io/portfolio/tct.html"),
    ("Portfolio", "puneetindeelhi940.github.io/portfolio/home.html"),
    ("Email", "puneet.ar@gmail.com"),
]
ly = Inches(5.3)
for label, url in links:
    add_text(sl, Inches(7.8), ly, Inches(1.5), Inches(0.25),
             label, font_size=8, color=FAINT, font_name='Consolas')
    add_text(sl, Inches(9.3), ly, Inches(3), Inches(0.25),
             url, font_size=8, color=WHITE, font_name='Consolas')
    ly += Inches(0.35)

add_text(sl, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.3),
         "Built with vibe coding  |  Double Diamond framework  |  June 2026",
         font_size=9, color=FAINT, font_name='Consolas', alignment=PP_ALIGN.CENTER)
add_page_num(sl, 30)


# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = '/home/user/portfolio/FeedIQ-TCT-Case-Study-Puneet-Arora.pptx'
prs.save(output_path)
print(f"OK saved to {output_path}")
print(f"Slides: {len(prs.slides)}")
