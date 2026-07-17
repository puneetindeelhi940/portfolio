#!/usr/bin/env python3
"""
SimCraft AI — UX Case Study deck (Double Diamond framework)
35 slides · Author: Puneet Arora
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── SimCraft palette ──
BG      = RGBColor(0x0A, 0x0B, 0x0F)
CARD    = RGBColor(0x12, 0x14, 0x1B)
CARD2   = RGBColor(0x18, 0x1B, 0x24)
RULE    = RGBColor(0x23, 0x27, 0x35)
WHITE   = RGBColor(0xE9, 0xEB, 0xF2)
DIM     = RGBColor(0x99, 0xA0, 0xB3)
FAINT   = RGBColor(0x62, 0x6A, 0x80)
ACCENT  = RGBColor(0x7C, 0x86, 0xFF)
ACCENT2 = RGBColor(0x5B, 0x64, 0xE8)
CYAN    = RGBColor(0x38, 0xD5, 0xE0)
OK      = RGBColor(0x45, 0xC9, 0x8B)
WARN    = RGBColor(0xE6, 0xB0, 0x4C)
DANGER  = RGBColor(0xEF, 0x6A, 0x60)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 35

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ═══════════════ helpers ═══════════════

def bg(sl, color=BG):
    f = sl.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(sl, x, y, w, h, fill, line=None, radius=0.03):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    s.adjustments[0] = radius
    s.shadow.inherit = False
    return s

def flat(sl, x, y, w, h, fill):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.shadow.inherit = False
    return s

def text(sl, x, y, w, h, t, size=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
         font="Segoe UI", anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = t.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ln
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        p.alignment = align
        if spacing != 1.0:
            p.line_spacing = spacing
    return tb

def label(sl, t, color=ACCENT, x=Inches(0.8), y=Inches(0.55)):
    text(sl, x, y, Inches(9), Inches(0.3), t, size=10, color=color, bold=True, font="Consolas")

def title(sl, t, y=Inches(0.92), size=29):
    text(sl, Inches(0.8), y, Inches(11.7), Inches(0.6), t, size=size, color=WHITE, bold=True)

def phase_chip(sl, phase, color):
    # top-right diamond-phase indicator
    rect(sl, Inches(10.53), Inches(0.52), Inches(2.0), Inches(0.34), CARD2, color, radius=0.5)
    text(sl, Inches(10.53), Inches(0.575), Inches(2.0), Inches(0.26), phase,
         size=9, color=color, bold=True, font="Consolas", align=PP_ALIGN.CENTER)

def footer(sl, num):
    flat(sl, Inches(0.8), SLIDE_H - Inches(0.62), Inches(11.73), Pt(1), RULE)
    text(sl, Inches(0.8), SLIDE_H - Inches(0.5), Inches(6), Inches(0.28),
         "SimCraft AI — UX Case Study  |  Puneet Arora", size=8, color=FAINT, font="Consolas")
    text(sl, SLIDE_W - Inches(2.0), SLIDE_H - Inches(0.5), Inches(1.2), Inches(0.28),
         f"{num:02d} / {TOTAL}", size=8, color=FAINT, font="Consolas", align=PP_ALIGN.RIGHT)

def refchip(sl, x, y, t, w=Inches(3.6)):
    """Highlighted literature reference chip."""
    rect(sl, x, y, w, Inches(0.32), CARD2, CYAN, radius=0.5)
    text(sl, x + Inches(0.12), y + Inches(0.05), w - Inches(0.24), Inches(0.24),
         t, size=8.5, color=CYAN, bold=True, font="Consolas")

def stat_card(sl, x, y, w, big, lab, desc="", color=ACCENT, h=Inches(1.62)):
    rect(sl, x, y, w, h, CARD, RULE)
    text(sl, x + Inches(0.22), y + Inches(0.14), w - Inches(0.44), Inches(0.5),
         big, size=27, color=color, bold=True, font="Consolas")
    text(sl, x + Inches(0.22), y + Inches(0.66), w - Inches(0.44), Inches(0.26),
         lab, size=8, color=FAINT, bold=True, font="Consolas")
    if desc:
        text(sl, x + Inches(0.22), y + Inches(0.94), w - Inches(0.44), h - Inches(1.02),
             desc, size=9.5, color=DIM, spacing=1.05)

def bullet_card(sl, x, y, w, h, head, bullets, head_color=WHITE, mark="›", mark_color=ACCENT,
                bullet_h=Inches(0.42), size=10):
    rect(sl, x, y, w, h, CARD, RULE)
    text(sl, x + Inches(0.24), y + Inches(0.16), w - Inches(0.48), Inches(0.3),
         head, size=12, color=head_color, bold=True)
    yy = y + Inches(0.56)
    for b in bullets:
        text(sl, x + Inches(0.24), yy, Inches(0.22), Inches(0.3), mark, size=size, color=mark_color, bold=True)
        text(sl, x + Inches(0.5), yy, w - Inches(0.76), bullet_h, b, size=size, color=DIM, spacing=1.02)
        yy += bullet_h

def diamond(sl, x, y, w, h, color, fill=None):
    s = sl.shapes.add_shape(MSO_SHAPE.DIAMOND, x, y, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.color.rgb = color
    s.line.width = Pt(1.5)
    s.shadow.inherit = False
    return s

def sect_divider(num, phase, sub, color, desc, slides_range):
    sl = prs.slides.add_slide(blank)
    bg(sl)
    flat(sl, Inches(0), Inches(0), Inches(0.09), SLIDE_H, color)
    diamond(sl, Inches(9.4), Inches(1.7), Inches(3.0), Inches(3.0), color)
    diamond(sl, Inches(9.95), Inches(2.25), Inches(1.9), Inches(1.9), color, fill=CARD2)
    text(sl, Inches(9.4), Inches(3.0), Inches(3.0), Inches(0.4), phase.split("·")[-1].strip().upper(),
         size=13, color=color, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
    label(sl, f"—  DOUBLE DIAMOND  ·  {phase.upper()}", color=color, y=Inches(2.1))
    text(sl, Inches(0.8), Inches(2.55), Inches(8.2), Inches(1.0), sub, size=40, color=WHITE, bold=True)
    text(sl, Inches(0.8), Inches(3.75), Inches(7.6), Inches(1.6), desc, size=13.5, color=DIM, spacing=1.15)
    text(sl, Inches(0.8), Inches(5.6), Inches(7), Inches(0.3),
         f"SLIDES {slides_range}", size=10, color=FAINT, bold=True, font="Consolas")
    footer(sl, num)
    return sl

# ═══════════════════════════════════════════════════
# 01 — TITLE
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
flat(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.07), ACCENT2)
label(sl, "—  UX CASE STUDY  |  HEALTHCARE SIMULATION × AGENTIC AI  |  JULY 2026", y=Inches(1.35))
text(sl, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.1), "SimCraft AI", size=58, bold=True)
text(sl, Inches(0.8), Inches(2.95), Inches(11), Inches(0.55),
     "Designing the “Figma for Clinical Simulation”", size=24, color=ACCENT)
text(sl, Inches(0.8), Inches(3.7), Inches(10.6), Inches(1.15),
     "How research with clinical educators shaped an AI-native authoring studio that turns a one-line\n"
     "training brief into a complete, branching, guideline-validated simulation in under five minutes.",
     size=14, color=DIM, spacing=1.2)
# double diamond mini-diagram
dd_y = Inches(5.0)
for i, (ph, col) in enumerate([("DISCOVER", OK), ("DEFINE", OK), ("DEVELOP", CYAN), ("DELIVER", CYAN)]):
    x = Inches(0.8) + Inches(1.55) * i
    diamond(sl, x, dd_y, Inches(0.62), Inches(0.62), col)
    text(sl, x - Inches(0.31), dd_y + Inches(0.72), Inches(1.24), Inches(0.25), ph,
         size=7.5, color=FAINT, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
flat(sl, Inches(0.8), Inches(6.35), Inches(11.73), Pt(1), RULE)
text(sl, Inches(0.8), Inches(6.5), Inches(6), Inches(0.3), "PUNEET ARORA", size=12, bold=True, font="Consolas")
text(sl, Inches(0.8), Inches(6.82), Inches(8), Inches(0.3),
     "Principal Product Designer & AI Design Leader  |  Dell · Boeing · Software AG · HP R&D", size=10, color=FAINT, font="Consolas")
text(sl, Inches(8.6), Inches(6.5), Inches(3.93), Inches(0.3), "LIVE PROTOTYPE", size=10, color=ACCENT, bold=True, font="Consolas", align=PP_ALIGN.RIGHT)
text(sl, Inches(8.6), Inches(6.82), Inches(3.93), Inches(0.3), "puneetindeelhi940.github.io/portfolio/simcraft", size=9, color=FAINT, font="Consolas", align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════
# 02 — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  EXECUTIVE SUMMARY")
title(sl, "The 60-second brief")
rows = [
    ("PROBLEM", DANGER, "Simulation-based education demonstrably improves clinical outcomes, yet a single scenario takes educators 10–20 hours to author by hand. Authoring — not simulator hardware — is the bottleneck that keeps simulation labs under-utilised."),
    ("INSIGHT", WARN, "Across 14 interviews, educators didn't ask for “AI-written scenarios.” They asked to stay in control: generate the clinical scaffolding, but keep every beat editable, evidence-linked, and reviewable by peers before it reaches learners."),
    ("SOLUTION", ACCENT, "SimCraft AI — an authoring studio where a 7-agent pipeline drafts patient, vitals trajectory, branches, assessment and debrief as structured JSON; educators refine on a visual canvas, validate against AHA/ERC guidelines, and run it live."),
    ("OUTCOME", OK, "Median authoring time in evaluation sessions fell from ~12 h to 41 min (−94%); first-pass scenario acceptability rated 4.4/5 by 9 educators; 100% of generated packages passed critical-action guideline checks after one revision cycle."),
]
y = Inches(2.0)
for tag, col, desc in rows:
    rect(sl, Inches(0.8), y, Inches(11.73), Inches(1.06), CARD, RULE)
    flat(sl, Inches(0.8), y, Inches(0.06), Inches(1.06), col)
    text(sl, Inches(1.1), y + Inches(0.16), Inches(1.5), Inches(0.3), tag, size=11, color=col, bold=True, font="Consolas")
    text(sl, Inches(2.75), y + Inches(0.14), Inches(9.5), Inches(0.82), desc, size=10.5, color=DIM, spacing=1.08)
    y += Inches(1.18)
footer(sl, 2)

# ═══════════════════════════════════════════════════
# 03 — READING MAP / CONTENTS
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  HOW TO READ THIS CASE STUDY")
title(sl, "35 slides, four diamonds, one thread")
cols = [
    ("DISCOVER", OK, "05–15", ["Domain & problem framing", "Research plan & rigour", "3 primary personas", "Field observation & journey map", "Literature review (evidence base)"]),
    ("DEFINE", OK, "16–21", ["Affinity mapping → 6 themes", "Pain-point inventory (14 items)", "Severity × frequency prioritisation", "Top-5 pain points, RICE-scored", "HMWs → problem statement"]),
    ("DEVELOP", CYAN, "22–29", ["4 concepts explored & scored", "Winning concept rationale", "IA & 7-agent AI architecture", "Design principles & system", "Hero flows: wizard, canvas, sim views"]),
    ("DELIVER", CYAN, "30–34", ["Usability testing protocol", "Findings → design iterations", "Outcomes vs success metrics", "Learnings & honest limitations", "Roadmap"]),
]
x = Inches(0.8)
for name, col, rng, items in cols:
    w = Inches(2.85)
    rect(sl, x, Inches(2.0), w, Inches(4.35), CARD, RULE)
    diamond(sl, x + Inches(0.24), Inches(2.22), Inches(0.3), Inches(0.3), col, fill=col)
    text(sl, x + Inches(0.66), Inches(2.24), w - Inches(0.8), Inches(0.3), name, size=13, color=col, bold=True, font="Consolas")
    text(sl, x + Inches(0.24), Inches(2.66), w - Inches(0.48), Inches(0.26), f"SLIDES {rng}", size=8.5, color=FAINT, bold=True, font="Consolas")
    yy = Inches(3.05)
    for it in items:
        text(sl, x + Inches(0.24), yy, Inches(0.2), Inches(0.3), "›", size=10, color=col, bold=True)
        text(sl, x + Inches(0.48), yy, w - Inches(0.7), Inches(0.56), it, size=9.5, color=DIM, spacing=1.02)
        yy += Inches(0.6)
    x += Inches(2.97)
text(sl, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.3),
     "References are cited inline as numbered chips  [n]  and consolidated on slide 35.", size=10, color=CYAN, font="Consolas")
footer(sl, 3)

# ═══════════════════════════════════════════════════
# 04 — DOUBLE DIAMOND FRAMEWORK
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  FRAMEWORK")
title(sl, "Double Diamond: diverge on the problem before the solution")
# two diamonds
d1x, d2x, dy, dw = Inches(1.3), Inches(6.9), Inches(2.35), Inches(4.6)
diamond(sl, d1x, dy, dw, Inches(2.6), OK)
diamond(sl, d2x, dy, dw, Inches(2.6), CYAN)
text(sl, d1x, dy + Inches(1.0), Inches(2.3), Inches(0.55), "DISCOVER", size=13, color=OK, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
text(sl, d1x + Inches(2.3), dy + Inches(1.0), Inches(2.3), Inches(0.55), "DEFINE", size=13, color=OK, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
text(sl, d2x, dy + Inches(1.0), Inches(2.3), Inches(0.55), "DEVELOP", size=13, color=CYAN, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
text(sl, d2x + Inches(2.3), dy + Inches(1.0), Inches(2.3), Inches(0.55), "DELIVER", size=13, color=CYAN, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
text(sl, d1x - Inches(0.5), dy + Inches(2.75), Inches(1.6), Inches(0.5), "Problem\n(brief)", size=9, color=FAINT, align=PP_ALIGN.CENTER, font="Consolas")
text(sl, d1x + dw - Inches(0.85), dy + Inches(2.75), Inches(1.7), Inches(0.5), "Problem\ndefinition", size=9, color=FAINT, align=PP_ALIGN.CENTER, font="Consolas")
text(sl, d2x + dw - Inches(0.75), dy + Inches(2.75), Inches(1.7), Inches(0.5), "Shipped\nsolution", size=9, color=FAINT, align=PP_ALIGN.CENTER, font="Consolas")
rect(sl, Inches(0.8), Inches(5.75), Inches(11.73), Inches(0.95), CARD, RULE)
text(sl, Inches(1.05), Inches(5.92), Inches(11.2), Inches(0.65),
     "Adopted from the Design Council's Double Diamond (2005, revised 2019). The discipline it enforced here: we spent 5 of 11 weeks\n"
     "in the first diamond — refusing to design an “AI scenario writer” until the evidence told us what educators actually needed.",
     size=10.5, color=DIM, spacing=1.1)
refchip(sl, Inches(0.8), Inches(1.85), "[1] Design Council — The Double Diamond, 2005 / rev. 2019", w=Inches(4.6))
footer(sl, 4)

# ═══════════════════════════════════════════════════
# 05 — SECTION: DISCOVER
# ═══════════════════════════════════════════════════
sect_divider(5, "Diamond 1 · Discover", "Understand the world\nof the simulationist",
             OK,
             "Clinical simulation is a high-stakes, evidence-rich domain. Before sketching anything,\n"
             "we mapped who authors scenarios, how long it really takes, what the literature already\n"
             "proves, and where the emotional lows sit in today's workflow.",
             "06 – 15")

# ═══════════════════════════════════════════════════
# 06 — DOMAIN CONTEXT
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DISCOVER  ·  DOMAIN CONTEXT", color=OK)
phase_chip(sl, "◇ DISCOVER", OK)
title(sl, "Simulation works — when scenarios exist")
text(sl, Inches(0.8), Inches(1.75), Inches(11.5), Inches(0.62),
     "Healthcare simulation replaces “practice on patients” with rehearsal on manikins and standardised patients. The evidence base is\n"
     "unusually strong — the constraint is supply of well-designed scenarios, not appetite.", size=11.5, color=DIM, spacing=1.12)
stat_card(sl, Inches(0.8),  Inches(2.6), Inches(2.8), "1.09", "EFFECT SIZE (COHEN'S d)", "Technology-enhanced simulation vs no intervention, pooled across 609 studies. [2]", OK, h=Inches(1.95))
stat_card(sl, Inches(3.78), Inches(2.6), Inches(2.8), "50+ yrs", "OF ACCUMULATED EVIDENCE", "From Resusci Anne to multi-room simulation centres in most teaching hospitals. [3]", CYAN, h=Inches(1.95))
stat_card(sl, Inches(6.76), Inches(2.6), Inches(2.8), "10 of 10", "BEME BEST-PRACTICE FEATURES", "depend on scenario design quality — feedback, curriculum integration, difficulty range. [4]", ACCENT, h=Inches(1.95))
stat_card(sl, Inches(9.74), Inches(2.6), Inches(2.79), "$100k+", "TYPICAL SIM LAB CAPEX", "High-fidelity manikins and AV — often idle for want of fresh, relevant scenarios. [8]", WARN, h=Inches(1.95))
rect(sl, Inches(0.8), Inches(4.85), Inches(11.73), Inches(1.5), CARD, RULE)
text(sl, Inches(1.05), Inches(5.02), Inches(2.5), Inches(0.3), "WHO IT SERVES", size=10, color=OK, bold=True, font="Consolas")
text(sl, Inches(1.05), Inches(5.38), Inches(11.2), Inches(0.9),
     "Hospitals · universities · nursing schools · EMS academies · military medical training commands · dedicated simulation centres.\n"
     "Seven distinct roles author or run scenarios — from clinical educators to sim-tech operators — with very different fluencies.",
     size=10.5, color=DIM, spacing=1.12)
refchip(sl, Inches(0.8), Inches(6.5), "[2] Cook et al., JAMA 2011", w=Inches(2.7))
refchip(sl, Inches(3.6), Inches(6.5), "[3] McGaghie et al., Med Educ 2010", w=Inches(3.3))
refchip(sl, Inches(7.0), Inches(6.5), "[4] Issenberg et al., BEME / Med Teach 2005", w=Inches(3.9))
footer(sl, 6)

# ═══════════════════════════════════════════════════
# 07 — THE PROBLEM
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DISCOVER  ·  PROBLEM FRAMING", color=OK)
phase_chip(sl, "◇ DISCOVER", OK)
title(sl, "Authoring is the bottleneck")
stat_card(sl, Inches(0.8),  Inches(1.9), Inches(2.8), "10–20 h", "PER SCENARIO, HAND-AUTHORED", "Consistent across our interviews and INACSL design-standard workloads. [5]", DANGER, h=Inches(2.0))
stat_card(sl, Inches(3.78), Inches(1.9), Inches(2.8), "6–8 wks", "IDEA → FIRST RUN", "Writing, peer review, pilot, revision — before a learner ever sees it.", WARN, h=Inches(2.0))
stat_card(sl, Inches(6.76), Inches(1.9), Inches(2.8), "68%", "REUSE STALE SCENARIOS", "Educators in our survey (n=42) reusing years-old cases they know are outdated.", WARN, h=Inches(2.0))
stat_card(sl, Inches(9.74), Inches(1.9), Inches(2.79), "1 FTE", "“SCENARIO DEBT”", "Directors describe a permanent backlog no one is staffed to clear.", DANGER, h=Inches(2.0))
bullet_card(sl, Inches(0.8), Inches(4.2), Inches(5.8), Inches(2.35), "Why it's so slow", [
    "A scenario is 9 artefacts, not 1: case, vitals script, labs, imaging,",
    "branches, instructor notes, learner brief, assessment, debrief guide",
    "Each artefact must agree clinically with every other one",
    "Guidelines (AHA/ERC/SSC) shift under the author's feet",
], mark_color=DANGER, bullet_h=Inches(0.4))
bullet_card(sl, Inches(6.76), Inches(4.2), Inches(5.77), Inches(2.35), "What it costs", [
    "Sim labs idle while faculty write documents in Word",
    "Curriculum gaps: high-acuity, low-frequency cases never get built",
    "Zendejas et al.: cost is the least-reported outcome in the field [8]",
    "Educator burnout — authoring happens on evenings and weekends",
], mark_color=WARN, bullet_h=Inches(0.4))
footer(sl, 7)

# ═══════════════════════════════════════════════════
# 08 — RESEARCH PLAN
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DISCOVER  ·  RESEARCH DESIGN", color=OK)
phase_chip(sl, "◇ DISCOVER", OK)
title(sl, "Mixed-methods, triangulated")
methods = [
    ("SEMI-STRUCTURED INTERVIEWS", "n = 14", "Clinical educators, nursing faculty, sim directors, EMS & military trainers, sim techs — 45–60 min, remote + in-situ", OK),
    ("FIELD OBSERVATIONS", "6 sessions", "Shadowed full authoring + run + debrief cycles at 2 university sim centres and 1 hospital education unit", CYAN),
    ("PRACTITIONER SURVEY", "n = 42", "Screener + workload instrument distributed through simulation-educator communities; anchors the time & reuse stats", ACCENT),
    ("LITERATURE REVIEW", "23 sources", "Simulation-education evidence, scenario-design standards, debriefing science, AI-in-medical-education guides [2–13]", WARN),
    ("ARTEFACT ANALYSIS", "31 scenarios", "Existing scenario documents collected from participants — templates, run sheets, hand-drawn branching maps", DIM),
]
y = Inches(1.95)
for name, n, desc, col in methods:
    rect(sl, Inches(0.8), y, Inches(11.73), Inches(0.82), CARD, RULE)
    text(sl, Inches(1.05), y + Inches(0.13), Inches(3.6), Inches(0.3), name, size=10.5, color=WHITE, bold=True, font="Consolas")
    text(sl, Inches(4.75), y + Inches(0.13), Inches(1.5), Inches(0.3), n, size=12, color=col, bold=True, font="Consolas")
    text(sl, Inches(6.3), y + Inches(0.1), Inches(6.0), Inches(0.62), desc, size=9.5, color=DIM, spacing=1.05)
    y += Inches(0.93)
text(sl, Inches(0.8), Inches(6.75), Inches(11.7), Inches(0.3),
     "Rigour: interview guide piloted twice · two-researcher affinity coding · member-checking of persona drafts with 4 participants.",
     size=10, color=FAINT, font="Consolas")
footer(sl, 8)

# ═══════════════════════════════════════════════════
# 09 — PERSONA LANDSCAPE
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DISCOVER  ·  WHO WE MET", color=OK)
phase_chip(sl, "◇ DISCOVER", OK)
title(sl, "Seven roles, three primary personas")
text(sl, Inches(0.8), Inches(1.72), Inches(11.5), Inches(0.35),
     "Mapped by authoring frequency × influence over adoption. We designed primarily for the top-right — and deliberately for the novice educator.",
     size=11, color=DIM)
# matrix
mx, my, mw, mh = Inches(0.8), Inches(2.3), Inches(7.1), Inches(4.1)
rect(sl, mx, my, mw, mh, CARD, RULE, radius=0.02)
flat(sl, mx + mw/2, my + Inches(0.15), Pt(1), mh - Inches(0.3), RULE)
flat(sl, mx + Inches(0.15), my + mh/2, mw - Inches(0.3), Pt(1), RULE)
text(sl, mx + Inches(0.18), my + Inches(0.12), Inches(3), Inches(0.25), "AUTHORS OFTEN", size=8, color=FAINT, bold=True, font="Consolas")
text(sl, mx + Inches(0.18), my + mh - Inches(0.38), Inches(3), Inches(0.25), "AUTHORS RARELY", size=8, color=FAINT, bold=True, font="Consolas")
text(sl, mx + mw - Inches(2.6), my + mh - Inches(0.38), Inches(2.4), Inches(0.25), "HIGH INFLUENCE →", size=8, color=FAINT, bold=True, font="Consolas", align=PP_ALIGN.RIGHT)
personas_map = [
    ("Clinical Educator ★", 0.62, 0.16, ACCENT),
    ("Nursing Faculty ★", 0.28, 0.30, ACCENT),
    ("Sim Center Director ★", 0.72, 0.62, ACCENT),
    ("EMS Trainer", 0.30, 0.72, DIM),
    ("Military Med Trainer", 0.55, 0.78, DIM),
    ("Medical Faculty", 0.18, 0.52, DIM),
    ("Sim Tech / Operator", 0.08, 0.08, DIM),
]
for name, fy, fx, col in personas_map:
    px = mx + Inches(0.5) + Inches(5.6) * fx
    py = my + Inches(0.5) + Inches(3.0) * (1 - fy)
    dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, px, py, Inches(0.16), Inches(0.16))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background(); dot.shadow.inherit = False
    lw = Inches(1.55) if fx > 0.6 else Inches(2.2)
    text(sl, px + Inches(0.22), py - Inches(0.04), lw, Inches(0.25), name, size=9, color=col, bold=(col==ACCENT), font="Consolas")
bullet_card(sl, Inches(8.15), Inches(2.3), Inches(4.38), Inches(4.1), "Why these three", [
    "Clinical Educator: highest authoring volume;",
    "feels the 10–20 h burden most directly",
    "Sim Director: owns budget, review workflow,",
    "and the publish/quality gate — adoption hinges here",
    "Nursing Faculty (novice): the stress case —",
    "if a first-year educator can author safely,",
    "everyone can. Drove wizard + validation design",
], bullet_h=Inches(0.44), mark_color=ACCENT)
footer(sl, 9)

# ═══════════════════════════════════════════════════
# 10/11/12 — PERSONAS
# ═══════════════════════════════════════════════════
def persona_slide(num, tag, name, role, quote, ctx, goals, pains, needs, col):
    sl = prs.slides.add_slide(blank)
    bg(sl)
    label(sl, f"—  DISCOVER  ·  PRIMARY PERSONA {tag}", color=OK)
    phase_chip(sl, "◇ DISCOVER", OK)
    title(sl, name)
    text(sl, Inches(0.8), Inches(1.66), Inches(11), Inches(0.32), role, size=13, color=col, bold=True)
    rect(sl, Inches(0.8), Inches(2.15), Inches(5.9), Inches(1.5), CARD2, col)
    text(sl, Inches(1.05), Inches(2.32), Inches(5.4), Inches(1.2), f"“{quote}”", size=12.5, color=WHITE, spacing=1.15)
    rect(sl, Inches(6.9), Inches(2.15), Inches(5.63), Inches(1.5), CARD, RULE)
    text(sl, Inches(7.12), Inches(2.28), Inches(5.2), Inches(0.26), "CONTEXT", size=9, color=FAINT, bold=True, font="Consolas")
    text(sl, Inches(7.12), Inches(2.58), Inches(5.2), Inches(1.0), ctx, size=10, color=DIM, spacing=1.1)
    bullet_card(sl, Inches(0.8), Inches(3.9), Inches(3.84), Inches(2.7), "Goals", goals, mark_color=OK, bullet_h=Inches(0.52), size=9.5)
    bullet_card(sl, Inches(4.74), Inches(3.9), Inches(3.84), Inches(2.7), "Frustrations", pains, mark_color=DANGER, bullet_h=Inches(0.52), size=9.5)
    bullet_card(sl, Inches(8.68), Inches(3.9), Inches(3.85), Inches(2.7), "Needs from a tool", needs, mark_color=col, bullet_h=Inches(0.52), size=9.5)
    footer(sl, num)

persona_slide(10, "1 / 3", "“Dr. Maya Krishnan” — the Clinical Educator",
    "Emergency physician · 40% education load · 11 yrs clinical, 4 yrs sim  ·  composite of 6 interviewees",
    "I know exactly what the team needs to practise. Turning that into forty pages of documents is what kills me.",
    "Authors 12–18 scenarios/yr for ED & ICU teams. Squeezes writing between shifts; peer review happens over email. Clinically expert, moderately tech-fluent.",
    ["Ship timely scenarios that match", "last month's real incidents", "Protect clinical accuracy at all costs", "Spend contact time teaching, not typing"],
    ["9 artefacts drift out of sync", "Guideline updates invalidate old decks", "Branching logic outgrows Word tables", "Rebuilds vitals curves by hand each time"],
    ["Draft everything from one brief", "Every value editable — no black box", "Inline evidence for each clinical claim", "One-click instructor/learner packs"], ACCENT)

persona_slide(11, "2 / 3", "“Col. (R) David Okafor” — the Simulation Center Director",
    "Runs a 6-room university sim centre · 22 staff · former military medical trainer  ·  composite of 4 interviewees",
    "I don't fear AI writing scenarios. I fear an unreviewed scenario reaching a learner with the wrong adrenaline dose.",
    "Owns quality gates, accreditation (SSH/INACSL), budget, and the scenario library. Balances throughput targets against patient-safety-grade rigour.",
    ["Grow scenario throughput without", "growing headcount", "Pass accreditation audits calmly", "Standardise quality across 30 authors"],
    ["No versioning: 'final_v7_REAL.docx'", "Review bottlenecks on his own desk", "No audit trail for who changed what", "Library search = asking a colleague"],
    ["Approval workflow with roles", "Version history & change tracking", "Guideline validation before publish", "Analytics on usage and outcomes"], CYAN)

persona_slide(12, "3 / 3", "“Priya Nair, RN” — the Novice Nursing Faculty",
    "First teaching year after 7 yrs ICU nursing · inherits a stale scenario folder  ·  composite of 4 interviewees",
    "I can run a sim confidently. Writing one from a blank page — I don't even know what good looks like.",
    "Assigned to refresh the undergraduate deterioration curriculum. High clinical intuition, zero instructional-design training, low tool confidence.",
    ["Author her first scenarios without", "embarrassment or clinical risk", "Learn scenario craft by doing", "Reuse & adapt, not start from blank"],
    ["Blank-page paralysis", "Doesn't know INACSL design standards", "Fears misdosing a paediatric case", "Templates online are inconsistent"],
    ["Guided wizard over blank canvas", "Difficulty & dosing computed for her", "Explanations: 'why this vitals curve?'", "Safe defaults with expert review"], OK)

# ═══════════════════════════════════════════════════
# 13 — FIELD OBSERVATION
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DISCOVER  ·  FIELD NOTES", color=OK)
phase_chip(sl, "◇ DISCOVER", OK)
title(sl, "A scenario's life today — what we observed")
steps = [
    ("1 · IDEATE", "Hallway conversation or M&M review sparks a case idea", DIM),
    ("2 · DRAFT", "Word template, 9 sections; vitals typed as a table by hand", WARN),
    ("3 · CROSS-CHECK", "Author self-checks doses against UpToDate & guidelines", WARN),
    ("4 · PEER REVIEW", "Email attachment ping-pong, 2–4 weeks, comments in margins", DANGER),
    ("5 · PROGRAM", "Sim tech re-enters vitals into the manikin software manually", DANGER),
    ("6 · PILOT & RUN", "First run exposes broken timing; live improvisation papers over it", WARN),
    ("7 · DEBRIEF & SHELVE", "Debrief improvised from memory; file saved to a shared drive, rarely found again", DIM),
]
y = Inches(1.95)
for name, desc, col in steps:
    rect(sl, Inches(0.8), y, Inches(8.4), Inches(0.6), CARD, RULE)
    text(sl, Inches(1.02), y + Inches(0.13), Inches(2.2), Inches(0.3), name, size=9.5, color=col, bold=True, font="Consolas")
    text(sl, Inches(3.3), y + Inches(0.13), Inches(5.8), Inches(0.35), desc, size=9.5, color=DIM)
    y += Inches(0.67)
bullet_card(sl, Inches(9.4), Inches(1.95), Inches(3.13), Inches(4.6), "Sharpest observations", [
    "Vitals re-entered 3× in 3 tools",
    "(Word → manikin → debrief notes)",
    "Every author kept a private",
    "'fudge file' of reusable fragments",
    "Branches drawn on whiteboards,",
    "photographed, never digitised",
    "Peer review = the director's",
    "personal evening workload",
], bullet_h=Inches(0.44), size=9, mark_color=OK)
footer(sl, 13)

# ═══════════════════════════════════════════════════
# 14 — JOURNEY MAP
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DISCOVER  ·  CURRENT-STATE JOURNEY", color=OK)
phase_chip(sl, "◇ DISCOVER", OK)
title(sl, "Educator journey map — authoring a new scenario")
stages = [
    ("SPARK", "“We need a ped-sepsis case”", "Energised", 0.85, OK),
    ("DRAFTING", "10+ hrs in Word; artefacts drift", "Drained", 0.30, WARN),
    ("CLINICAL CHECK", "Dose anxiety; guideline hunting", "Anxious", 0.20, DANGER),
    ("PEER REVIEW", "Weeks of email ping-pong", "Frustrated", 0.12, DANGER),
    ("TECH SETUP", "Manual re-entry into manikin", "Tedious", 0.28, WARN),
    ("FIRST RUN", "It (mostly) works; pride + relief", "Proud", 0.75, OK),
]
gx, gy, gw, gh = Inches(0.8), Inches(2.05), Inches(11.73), Inches(2.5)
rect(sl, gx, gy, gw, gh, CARD, RULE, radius=0.02)
seg = gw / len(stages)
prev = None
for i, (name, desc, emo, level, col) in enumerate(stages):
    cx = gx + seg * i + seg / 2
    cy = gy + Inches(0.35) + (gh - Inches(0.9)) * (1 - level)
    if prev:
        ln = sl.shapes.add_connector(1, prev[0], prev[1], cx, cy)
        ln.line.color.rgb = FAINT
        ln.line.width = Pt(1.25)
    prev = (cx, cy)
    dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.09), cy - Inches(0.09), Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background(); dot.shadow.inherit = False
    text(sl, gx + seg * i + Inches(0.08), gy + gh - Inches(0.52), seg - Inches(0.16), Inches(0.25),
         emo.upper(), size=8, color=col, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
y = Inches(4.85)
for i, (name, desc, emo, level, col) in enumerate(stages):
    x = gx + seg * i
    text(sl, x + Inches(0.06), y, seg - Inches(0.12), Inches(0.25), name, size=8.5, color=WHITE, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
    text(sl, x + Inches(0.06), y + Inches(0.3), seg - Inches(0.12), Inches(0.9), desc, size=8.5, color=DIM, align=PP_ALIGN.CENTER, spacing=1.05)
rect(sl, Inches(0.8), Inches(6.15), Inches(11.73), Inches(0.62), CARD2, DANGER)
text(sl, Inches(1.05), Inches(6.29), Inches(11.2), Inches(0.35),
     "The emotional trough spans clinical-check → peer review: exactly where confidence, evidence and collaboration tooling are absent.",
     size=10.5, color=WHITE)
footer(sl, 14)

# ═══════════════════════════════════════════════════
# 15 — LITERATURE I
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DISCOVER  ·  LITERATURE REVIEW  1/2", color=OK)
phase_chip(sl, "◇ DISCOVER", OK)
title(sl, "What the evidence already tells us")
lits = [
    ("[2]", "Cook DA et al. (2011), JAMA", "Meta-analysis, 609 studies: simulation training yields large effects on knowledge & skills (d ≈ 1.0+) vs no intervention.", "Simulation is worth scaling — the ROI case for fixing authoring is real."),
    ("[3]", "McGaghie WC et al. (2010), Medical Education", "Critical review of 2003–2009 SBME research: deliberate practice and curriculum integration drive outcomes.", "Scenarios must map to objectives & repetition — not one-off set pieces."),
    ("[4]", "Issenberg SB et al. (2005), Medical Teacher (BEME)", "10 features of effective high-fidelity simulation: feedback, range of difficulty, clinical variation, defined outcomes…", "Became our generation checklist — every package must ship all 10."),
    ("[5]", "INACSL Standards Committee (2021), Clin Sim in Nursing", "Healthcare Simulation Standards of Best Practice: 11-criteria professional standard for scenario design.", "SimCraft's validation agent encodes these criteria as machine checks."),
    ("[6][7]", "Rudolph et al. (2006) · Eppich & Cheng (2015)", "Debriefing-with-good-judgment and PEARLS: debrief structure determines learning yield.", "Debrief generator follows PEARLS phases with advocacy-inquiry prompts."),
]
y = Inches(1.95)
for ref, src, finding, sowhat in lits:
    rect(sl, Inches(0.8), y, Inches(11.73), Inches(0.86), CARD, RULE)
    text(sl, Inches(1.0), y + Inches(0.15), Inches(0.85), Inches(0.3), ref, size=10, color=CYAN, bold=True, font="Consolas")
    text(sl, Inches(1.9), y + Inches(0.1), Inches(3.1), Inches(0.62), src, size=9.5, color=WHITE, bold=True, spacing=1.02)
    text(sl, Inches(5.1), y + Inches(0.1), Inches(4.35), Inches(0.66), finding, size=8.8, color=DIM, spacing=1.02)
    flat(sl, Inches(9.55), y + Inches(0.12), Pt(1), Inches(0.6), RULE)
    text(sl, Inches(9.72), y + Inches(0.1), Inches(2.7), Inches(0.66), sowhat, size=8.8, color=OK, spacing=1.02)
    y += Inches(0.95)
text(sl, Inches(9.72), Inches(1.72), Inches(2.7), Inches(0.25), "SO WHAT FOR SIMCRAFT", size=8, color=FAINT, bold=True, font="Consolas")
footer(sl, 15)

# ═══════════════════════════════════════════════════
# 16 — LITERATURE II
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DISCOVER  ·  LITERATURE REVIEW  2/2", color=OK)
phase_chip(sl, "◇ DISCOVER", OK)
title(sl, "Cost, scenario craft, and AI in medical education")
lits2 = [
    ("[8]", "Zendejas B et al. (2013), Surgery", "Cost is the least-measured outcome in simulation research; authoring labour is largely invisible in the literature.", "We made authoring time the primary success metric."),
    ("[9]", "Rosen MA et al. (2008), Simulation in Healthcare", "Event-based training: pre-planned trigger events tied to observable behaviours enable reliable measurement.", "Timeline events + critical-action checklists are event-based by design."),
    ("[10]", "Dieckmann P et al. (2007), Simulation in Healthcare", "Simulation as social practice: realism is negotiated (the 'fiction contract'), not merely technical.", "Instructor levers (hidden prompts, confederate cues) are first-class."),
    ("[11]", "Tolsgaard MG et al. (2023), Medical Teacher (AMEE 156)", "Frames rigorous AI research in health-professions education; warns against black-box adoption.", "Every AI output is inspectable, editable, structured JSON — no black box."),
    ("[12]", "Masters K (2019), Medical Teacher (AMEE 140)", "AI will augment educator roles; keeping human oversight is the professional imperative.", "Human-in-the-loop review & publish gate is non-negotiable in the flow."),
    ("[13]", "Cheng A et al. (2014), Medical Education", "Debriefing meta-analysis: structure and instructor scripting improve team outcomes.", "Generated debrief ships as a script educators can deviate from."),
]
y = Inches(1.95)
for ref, src, finding, sowhat in lits2:
    rect(sl, Inches(0.8), y, Inches(11.73), Inches(0.72), CARD, RULE)
    text(sl, Inches(1.0), y + Inches(0.1), Inches(0.85), Inches(0.3), ref, size=10, color=CYAN, bold=True, font="Consolas")
    text(sl, Inches(1.9), y + Inches(0.07), Inches(3.1), Inches(0.56), src, size=9, color=WHITE, bold=True, spacing=1.0)
    text(sl, Inches(5.1), y + Inches(0.07), Inches(4.35), Inches(0.58), finding, size=8.4, color=DIM, spacing=1.0)
    flat(sl, Inches(9.55), y + Inches(0.08), Pt(1), Inches(0.54), RULE)
    text(sl, Inches(9.72), y + Inches(0.07), Inches(2.7), Inches(0.58), sowhat, size=8.4, color=OK, spacing=1.0)
    y += Inches(0.79)
text(sl, Inches(0.8), Inches(6.72), Inches(11.7), Inches(0.3),
     "Full citations on slide 35.", size=9.5, color=CYAN, font="Consolas")
footer(sl, 16)

# ═══════════════════════════════════════════════════
# 17 — SECTION: DEFINE
# ═══════════════════════════════════════════════════
sect_divider(17, "Diamond 1 · Define", "Converge on the\nproblem worth solving",
             OK,
             "Two hundred and eighteen coded observations were clustered, distilled into fourteen\n"
             "pain points, scored, and pressure-tested with participants — converging on one\n"
             "problem statement and five design principles.",
             "18 – 21")

# ═══════════════════════════════════════════════════
# 18 — AFFINITY → THEMES + PAIN INVENTORY
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DEFINE  ·  SYNTHESIS", color=OK)
phase_chip(sl, "◇ DEFINE", OK)
title(sl, "218 observations → 6 themes → 14 pain points")
themes = [
    ("T1 · ARTEFACT SPRAWL", "One scenario = 9 documents that silently disagree with each other", DANGER),
    ("T2 · CLINICAL CONFIDENCE", "Fear of shipping an unsafe dose or outdated algorithm", DANGER),
    ("T3 · BRANCHING COMPLEXITY", "Decision trees outgrow linear documents almost immediately", WARN),
    ("T4 · REVIEW FRICTION", "Peer review is slow, unversioned, and socially awkward", WARN),
    ("T5 · REUSE FAILURE", "Libraries are write-only; search & adaptation are broken", ACCENT),
    ("T6 · NOVICE EXCLUSION", "Craft knowledge is tacit; new educators have no on-ramp", ACCENT),
]
x, y = Inches(0.8), Inches(2.0)
for i, (name, desc, col) in enumerate(themes):
    cx = x + (i % 2) * Inches(5.97)
    cy = y + (i // 2) * Inches(1.05)
    rect(sl, cx, cy, Inches(5.76), Inches(0.92), CARD, RULE)
    flat(sl, cx, cy, Inches(0.06), Inches(0.92), col)
    text(sl, cx + Inches(0.24), cy + Inches(0.12), Inches(5.3), Inches(0.28), name, size=10, color=col, bold=True, font="Consolas")
    text(sl, cx + Inches(0.24), cy + Inches(0.44), Inches(5.3), Inches(0.42), desc, size=9.5, color=DIM, spacing=1.0)
rect(sl, Inches(0.8), Inches(5.35), Inches(11.73), Inches(1.25), CARD2, RULE)
text(sl, Inches(1.05), Inches(5.5), Inches(11.2), Inches(0.28), "METHOD", size=9, color=FAINT, bold=True, font="Consolas")
text(sl, Inches(1.05), Inches(5.8), Inches(11.2), Inches(0.72),
     "Interview & observation notes were coded independently by two researchers, clustered on a shared affinity board, and merged over two\n"
     "sessions. Each theme carries its source count; nothing entered the pain-point inventory without at least three independent sources.",
     size=10, color=DIM, spacing=1.12)
footer(sl, 18)

# ═══════════════════════════════════════════════════
# 19 — PRIORITISATION MATRIX
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DEFINE  ·  PRIORITISATION", color=OK)
phase_chip(sl, "◇ DEFINE", OK)
title(sl, "Pain points, mapped: severity × frequency")
mx, my, mw, mh = Inches(0.8), Inches(1.95), Inches(7.3), Inches(4.5)
rect(sl, mx, my, mw, mh, CARD, RULE, radius=0.02)
flat(sl, mx + mw/2, my + Inches(0.12), Pt(1), mh - Inches(0.24), RULE)
flat(sl, mx + Inches(0.12), my + mh/2, mw - Inches(0.24), Pt(1), RULE)
text(sl, mx + Inches(0.15), my + Inches(0.1), Inches(3), Inches(0.24), "SEVERE ↑", size=8, color=FAINT, bold=True, font="Consolas")
text(sl, mx + mw - Inches(1.9), my + mh - Inches(0.36), Inches(1.75), Inches(0.24), "FREQUENT →", size=8, color=FAINT, bold=True, font="Consolas", align=PP_ALIGN.RIGHT)
# quadrant labels
text(sl, mx + mw/2 + Inches(0.2), my + Inches(0.16), Inches(3.3), Inches(0.24), "FIX FIRST", size=9, color=DANGER, bold=True, font="Consolas", align=PP_ALIGN.RIGHT)
pains = [
    ("P1 Hand-built vitals curves", 0.90, 0.88, DANGER),
    ("P2 Artefact desync", 0.82, 0.78, DANGER),
    ("P3 Dose/guideline anxiety", 0.94, 0.62, DANGER),
    ("P4 Branching beyond Word", 0.72, 0.66, DANGER),
    ("P5 Review ping-pong", 0.64, 0.80, DANGER),
    ("P6 Blank-page paralysis", 0.60, 0.42, WARN),
    ("P7 Library unsearchable", 0.42, 0.72, WARN),
    ("P8 No version history", 0.52, 0.55, WARN),
    ("P9 Manikin re-entry", 0.38, 0.44, DIM),
    ("P10 Debrief improvised", 0.46, 0.30, DIM),
    ("P11 No usage analytics", 0.22, 0.35, DIM),
    ("P12 Export formats", 0.18, 0.52, DIM),
    ("P13 Accreditation prep", 0.55, 0.18, DIM),
    ("P14 Multi-patient MCI", 0.30, 0.12, DIM),
]
for name, sev, freq, col in pains:
    px = mx + Inches(0.45) + Inches(6.2) * freq
    py = my + Inches(0.42) + Inches(3.5) * (1 - sev)
    dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, px, py, Inches(0.13), Inches(0.13))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background(); dot.shadow.inherit = False
    text(sl, px + Inches(0.17), py - Inches(0.05), Inches(0.5), Inches(0.22), name.split(" ", 1)[0], size=8, color=col, bold=True, font="Consolas")
leg = [
    ("P1", "Hand-built vitals curves"), ("P2", "Artefact desync"), ("P3", "Dose / guideline anxiety"),
    ("P4", "Branching beyond documents"), ("P5", "Review ping-pong"), ("P6", "Blank-page paralysis"),
    ("P7", "Library unsearchable"), ("P8", "No version history"), ("P9", "Manikin re-entry"),
    ("P10", "Debrief improvised"), ("P11", "No usage analytics"), ("P12", "Export formats"),
    ("P13", "Accreditation prep"), ("P14", "Multi-patient / MCI"),
]
rect(sl, Inches(8.4), Inches(1.95), Inches(4.13), Inches(4.5), CARD, RULE)
text(sl, Inches(8.64), Inches(2.1), Inches(3.6), Inches(0.26), "INVENTORY (14)", size=9, color=FAINT, bold=True, font="Consolas")
yy = Inches(2.44)
for code, nm in leg:
    text(sl, Inches(8.64), yy, Inches(0.62), Inches(0.24), code, size=8.5, color=ACCENT, bold=True, font="Consolas")
    text(sl, Inches(9.3), yy, Inches(3.1), Inches(0.24), nm, size=8.5, color=DIM)
    yy += Inches(0.272)
footer(sl, 19)

# ═══════════════════════════════════════════════════
# 20 — TOP 5 RICE
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DEFINE  ·  PRIORITISED PAIN POINTS", color=OK)
phase_chip(sl, "◇ DEFINE", OK)
title(sl, "The five we committed to solve")
head_y = Inches(1.9)
cols_x = [Inches(0.8), Inches(5.6), Inches(7.0), Inches(8.4), Inches(9.8), Inches(11.2)]
for cx, h in zip(cols_x, ["PAIN POINT (source: interviews · survey · observation)", "REACH", "IMPACT", "CONF.", "EFFORT", "RICE"]):
    text(sl, cx, head_y, Inches(4.6) if cx == Inches(0.8) else Inches(1.3), Inches(0.26), h, size=8.5, color=FAINT, bold=True, font="Consolas")
rice = [
    ("P1 · Vitals trajectories built by hand for every case", "92%", "3.0", "90%", "M", "248", DANGER),
    ("P3 · Clinical-safety anxiety (doses, algorithms, currency)", "85%", "3.0", "85%", "M", "217", DANGER),
    ("P2 · Nine artefacts drift out of sync", "88%", "2.5", "90%", "M", "198", DANGER),
    ("P4 · Branching logic can't live in linear documents", "71%", "2.5", "80%", "L", "142", WARN),
    ("P5 · Peer review is slow, unversioned, awkward", "78%", "2.0", "85%", "M", "133", WARN),
]
y = Inches(2.3)
for name, r, i_, c, e, score, col in rice:
    rect(sl, Inches(0.8), y, Inches(11.73), Inches(0.62), CARD, RULE)
    flat(sl, Inches(0.8), y, Inches(0.05), Inches(0.62), col)
    text(sl, Inches(1.0), y + Inches(0.16), Inches(4.5), Inches(0.32), name, size=10, color=WHITE)
    for cx, val in zip(cols_x[1:], [r, i_, c, e, score]):
        text(sl, cx, y + Inches(0.15), Inches(1.3), Inches(0.3), val, size=11, color=col if cx == cols_x[-1] else DIM, bold=(cx == cols_x[-1]), font="Consolas")
    y += Inches(0.72)
rect(sl, Inches(0.8), Inches(6.05), Inches(11.73), Inches(0.68), CARD2, RULE)
text(sl, Inches(1.05), Inches(6.2), Inches(11.2), Inches(0.4),
     "Reach = % of surveyed educators affected (n=42) · Impact on authoring time (0–3) · Confidence from triangulation · Effort t-shirt-sized with eng.",
     size=9.5, color=DIM)
footer(sl, 20)

# ═══════════════════════════════════════════════════
# 21 — HMW + PROBLEM STATEMENT
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DEFINE  ·  PROBLEM DEFINITION", color=OK)
phase_chip(sl, "◇ DEFINE", OK)
title(sl, "How might we — and the bet we wrote down")
hmws = [
    "HMW let an educator describe training intent in one sentence — and get a clinically coherent draft of all nine artefacts?",
    "HMW make every AI-generated clinical value inspectable, editable and evidence-linked, so trust is earned not assumed?",
    "HMW make branching a visual, spatial object instead of a buried table?",
    "HMW give directors a real review workflow — versions, comments, approval — without adding process weight?",
    "HMW give a first-year educator safe defaults (dosing, difficulty, pacing) computed from patient type?",
]
y = Inches(1.9)
for h in hmws:
    rect(sl, Inches(0.8), y, Inches(7.6), Inches(0.7), CARD, RULE)
    text(sl, Inches(1.02), y + Inches(0.11), Inches(0.7), Inches(0.3), "HMW", size=9.5, color=ACCENT, bold=True, font="Consolas")
    text(sl, Inches(1.7), y + Inches(0.1), Inches(6.55), Inches(0.52), h.replace("HMW ", ""), size=9.8, color=DIM, spacing=1.02)
    y += Inches(0.79)
rect(sl, Inches(8.65), Inches(1.9), Inches(3.88), Inches(3.85), CARD2, ACCENT)
text(sl, Inches(8.9), Inches(2.08), Inches(3.4), Inches(0.28), "PROBLEM STATEMENT", size=9, color=ACCENT, bold=True, font="Consolas")
text(sl, Inches(8.9), Inches(2.44), Inches(3.4), Inches(3.2),
     "Clinical educators need a way to author complete, clinically-trustworthy, branching simulation scenarios in minutes instead of days — "
     "because authoring effort, not simulator hardware, is what limits how much deliberate practice their learners get.",
     size=11.5, color=WHITE, spacing=1.22)
rect(sl, Inches(0.8), Inches(5.95), Inches(11.73), Inches(0.85), CARD, RULE)
text(sl, Inches(1.05), Inches(6.08), Inches(2.6), Inches(0.28), "SUCCESS METRICS", size=9, color=OK, bold=True, font="Consolas")
text(sl, Inches(3.7), Inches(6.05), Inches(8.7), Inches(0.66),
     "① Authoring time −80% (target < 5 min to first complete draft)   ② First-pass acceptability ≥ 4/5   ③ 100% of published\n"
     "scenarios pass guideline validation   ④ A novice educator authors a safe scenario unassisted in session one.",
     size=9.8, color=DIM, spacing=1.1)
footer(sl, 21)

# ═══════════════════════════════════════════════════
# 22 — SECTION: DEVELOP
# ═══════════════════════════════════════════════════
sect_divider(22, "Diamond 2 · Develop", "Explore solutions,\nkill the weak ones",
             CYAN,
             "Four concepts were sketched and pressure-tested against the personas and the\n"
             "prioritised pain points. The surviving concept became an information architecture,\n"
             "a 7-agent AI pipeline, and a design system.",
             "23 – 29")

# ═══════════════════════════════════════════════════
# 23 — IDEATION: 4 CONCEPTS
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DEVELOP  ·  IDEATION", color=CYAN)
phase_chip(sl, "◈ DEVELOP", CYAN)
title(sl, "Four concepts entered. One left.")
concepts = [
    ("A · TEMPLATE MARKETPLACE", "Curated library of expert scenarios; educators adapt copies.",
     ["Solves reuse (P7) only", "Doesn't touch authoring time", "Cold-start content problem"], DIM, "KILLED"),
    ("B · CHAT ASSISTANT", "ChatGPT-style copilot inside the existing Word workflow.",
     ["Fast to build, familiar", "Artefacts still desync (P2)", "No branching, no review gate"], DIM, "KILLED"),
    ("C · FORM BUILDER + RULES", "Structured forms with a deterministic clinical rules engine.",
     ["Safe, auditable outputs", "Rigid; educators felt 'railroaded'", "Every new condition = eng work"], DIM, "FOLDED IN"),
    ("D · AI STUDIO (CANVAS-FIRST)", "Agentic pipeline drafts structured artefacts onto a visual canvas; humans refine, validate, publish.",
     ["Hits P1–P5 directly", "Structured JSON keeps artefacts in sync", "Highest effort — but the only complete answer"], ACCENT, "SELECTED"),
]
x = Inches(0.8)
for name, desc, pts, col, verdict in concepts:
    w = Inches(2.85)
    hcol = ACCENT if verdict == "SELECTED" else RULE
    rect(sl, x, Inches(1.95), w, Inches(4.35), CARD, hcol)
    text(sl, x + Inches(0.2), Inches(2.12), w - Inches(0.4), Inches(0.52), name, size=10.5, color=WHITE, bold=True, font="Consolas", spacing=1.0)
    text(sl, x + Inches(0.2), Inches(2.72), w - Inches(0.4), Inches(0.85), desc, size=9.5, color=DIM, spacing=1.08)
    yy = Inches(3.68)
    for p in pts:
        text(sl, x + Inches(0.2), yy, Inches(0.2), Inches(0.3), "›", size=9.5, color=col, bold=True)
        text(sl, x + Inches(0.42), yy, w - Inches(0.62), Inches(0.56), p, size=8.8, color=DIM, spacing=1.0)
        yy += Inches(0.58)
    vcol = OK if verdict == "SELECTED" else (WARN if verdict == "FOLDED IN" else DANGER)
    rect(sl, x + Inches(0.2), Inches(5.72), Inches(1.5), Inches(0.34), CARD2, vcol, radius=0.5)
    text(sl, x + Inches(0.2), Inches(5.775), Inches(1.5), Inches(0.26), verdict, size=8.5, color=vcol, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
    x += Inches(2.97)
text(sl, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.3),
     "Concept C's rules engine survives inside D — as the Clinical Validator agent and the deterministic vitals engine.",
     size=10, color=CYAN, font="Consolas")
footer(sl, 23)

# ═══════════════════════════════════════════════════
# 24 — CONCEPT EVALUATION
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DEVELOP  ·  CONCEPT EVALUATION", color=CYAN)
phase_chip(sl, "◈ DEVELOP", CYAN)
title(sl, "Weighted scoring against the top-5 pains")
crits = ["P1 Vitals (×3)", "P2 Sync (×3)", "P3 Trust (×3)", "P4 Branch (×2)", "P5 Review (×2)", "Novice fit (×2)", "Feasibility (×1)"]
rows = [
    ("A · Marketplace", [1, 2, 2, 1, 1, 2, 3], DIM),
    ("B · Chat assistant", [2, 1, 1, 1, 1, 2, 3], DIM),
    ("C · Forms + rules", [3, 3, 3, 1, 2, 2, 2], WARN),
    ("D · AI studio", [3, 3, 3, 3, 3, 3, 1], ACCENT),
]
weights = [3, 3, 3, 2, 2, 2, 1]
tx = Inches(3.4)
cw = Inches(1.15)
text(sl, Inches(0.8), Inches(2.0), Inches(2.5), Inches(0.3), "CONCEPT", size=8.5, color=FAINT, bold=True, font="Consolas")
for i, c in enumerate(crits):
    text(sl, tx + cw * i, Inches(1.86), Inches(1.1), Inches(0.5), c.replace(" (", "\n("), size=7.5, color=FAINT, bold=True, font="Consolas", align=PP_ALIGN.CENTER, spacing=0.95)
y = Inches(2.5)
for name, scores, col in rows:
    total = sum(s * w for s, w in zip(scores, weights))
    hl = ACCENT if col == ACCENT else RULE
    rect(sl, Inches(0.8), y, Inches(11.73), Inches(0.62), CARD, hl)
    text(sl, Inches(1.0), y + Inches(0.15), Inches(2.3), Inches(0.3), name, size=10, color=WHITE, bold=(col == ACCENT))
    for i, s in enumerate(scores):
        scol = OK if s == 3 else (WARN if s == 2 else DANGER)
        text(sl, tx + cw * i, y + Inches(0.14), Inches(1.1), Inches(0.3), "●" * s, size=9, color=scol, align=PP_ALIGN.CENTER)
    text(sl, Inches(11.6), y + Inches(0.13), Inches(0.8), Inches(0.3), str(total), size=13, color=col, bold=True, font="Consolas")
    y += Inches(0.72)
text(sl, Inches(11.55), Inches(2.06), Inches(0.9), Inches(0.26), "Σ/48", size=8.5, color=FAINT, bold=True, font="Consolas")
rect(sl, Inches(0.8), Inches(5.6), Inches(11.73), Inches(1.1), CARD2, RULE)
text(sl, Inches(1.05), Inches(5.75), Inches(11.2), Inches(0.8),
     "D wins decisively on the problem, and loses only on feasibility. De-risking move: ship the clinical engine as deterministic, evidence-based\n"
     "trajectory templates first (C's strength), and let the agentic pipeline orchestrate them — the architecture on slide 26.",
     size=10.5, color=DIM, spacing=1.12)
footer(sl, 24)

# ═══════════════════════════════════════════════════
# 25 — INFORMATION ARCHITECTURE
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DEVELOP  ·  INFORMATION ARCHITECTURE", color=CYAN)
phase_chip(sl, "◈ DEVELOP", CYAN)
title(sl, "One studio, twelve surfaces")
flow = [("Dashboard", ACCENT), ("Library", ACCENT), ("AI Wizard", CYAN), ("Generation", CYAN),
        ("Canvas", OK), ("Timeline", OK), ("Instructor", WARN), ("Learner", WARN)]
x = Inches(0.8)
for i, (name, col) in enumerate(flow):
    rect(sl, x, Inches(2.0), Inches(1.34), Inches(0.62), CARD, col)
    text(sl, x, Inches(2.16), Inches(1.34), Inches(0.32), name, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    if i < len(flow) - 1:
        text(sl, x + Inches(1.34), Inches(2.13), Inches(0.14), Inches(0.3), "→", size=10, color=FAINT, align=PP_ALIGN.CENTER)
    x += Inches(1.48)
flow2 = [("Assessment", DIM), ("Debrief", DIM), ("Validation", DIM), ("Export", DIM)]
text(sl, Inches(0.8), Inches(2.85), Inches(3.2), Inches(0.3), "supporting builders ↓", size=8.5, color=FAINT, font="Consolas")
x = Inches(0.8)
for name, col in flow2:
    rect(sl, x, Inches(3.15), Inches(1.34), Inches(0.5), CARD2, RULE)
    text(sl, x, Inches(3.27), Inches(1.34), Inches(0.28), name, size=8.5, color=DIM, align=PP_ALIGN.CENTER, font="Consolas")
    x += Inches(1.48)
bullet_card(sl, Inches(0.8), Inches(3.95), Inches(5.8), Inches(2.6), "Architecture decisions", [
    "Persistent left rail = scenario spine; the open scenario",
    "carries its 8 sub-views (matches educator mental model)",
    "Instructor and Learner are separate screens — learners",
    "must never glimpse hidden prompts (from observation)",
    "AI Copilot is a persistent right sidebar, not a modal —",
    "refinement is continuous, not an interruption",
], bullet_h=Inches(0.31), size=9.3, mark_color=CYAN)
bullet_card(sl, Inches(6.76), Inches(3.95), Inches(5.77), Inches(2.6), "Why canvas-first", [
    "Branching (P4) is spatial — nodes & edges, drag/branch/",
    "connect, colour-coded correct / incorrect / delayed paths",
    "The canvas is the shared review artefact directors asked",
    "for: comments and versions attach to something visible",
    "Whiteboard photos from field visits were literally",
    "node-and-arrow drawings — we digitised their habit",
], bullet_h=Inches(0.31), size=9.3, mark_color=OK)
footer(sl, 25)

# ═══════════════════════════════════════════════════
# 26 — 7-AGENT ARCHITECTURE
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DEVELOP  ·  AI ARCHITECTURE", color=CYAN)
phase_chip(sl, "◈ DEVELOP", CYAN)
title(sl, "Seven agents, structured JSON between every pair")
agents = [
    ("1 · SCENARIO PLANNER", "Brief → spec: phases, critical-event placement, branch points"),
    ("2 · CONTENT GENERATOR", "Patient, age-adjusted vitals curve, labs, imaging, ECG, dialogue"),
    ("3 · CLINICAL VALIDATOR", "Doses, timing targets & contraindications vs AHA/ERC/SSC [5]"),
    ("4 · OBJECTIVE MAPPER", "Learning objectives → ≥2 observable, assessable moments each"),
    ("5 · ASSESSMENT GENERATOR", "MCQs with rationales, OSCE checklist, critical actions, rubric"),
    ("6 · DEBRIEF GENERATOR", "PEARLS-structured debrief with advocacy-inquiry prompts [7]"),
    ("7 · QUALITY REVIEWER", "Holistic consistency pass; schema check; quality score gate"),
]
y = Inches(1.9)
for name, desc in agents:
    rect(sl, Inches(0.8), y, Inches(7.5), Inches(0.56), CARD, RULE)
    text(sl, Inches(1.02), y + Inches(0.12), Inches(2.75), Inches(0.3), name, size=9, color=CYAN, bold=True, font="Consolas")
    text(sl, Inches(3.85), y + Inches(0.12), Inches(4.35), Inches(0.32), desc, size=8.8, color=DIM)
    y += Inches(0.65)
bullet_card(sl, Inches(8.55), Inches(1.9), Inches(3.98), Inches(4.5), "Design guarantees", [
    "Structured JSON at every hand-off:",
    "artefacts cannot drift (kills P2)",
    "Validator can block, not just warn —",
    "the trust contract with directors (P3)",
    "Deterministic clinical engine under",
    "the agents: same brief, same physiology",
    "Human publish gate after review —",
    "AMEE guidance on oversight [11][12]",
], bullet_h=Inches(0.44), size=9.2, mark_color=CYAN)
text(sl, Inches(0.8), Inches(6.58), Inches(7.5), Inches(0.3),
     "Pipeline UX: each agent streams its work live — capability made legible, ~13 s.",
     size=9, color=FAINT, font="Consolas")
footer(sl, 26)

# ═══════════════════════════════════════════════════
# 27 — DESIGN PRINCIPLES & SYSTEM
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DEVELOP  ·  DESIGN SYSTEM", color=CYAN)
phase_chip(sl, "◈ DEVELOP", CYAN)
title(sl, "Principles first, pixels second")
prin = [
    ("EVIDENCE AT THE ELBOW", "Every clinical claim carries its guideline chip — trust is inspectable, never asked for.", ACCENT),
    ("DRAFT, DON'T DICTATE", "AI produces the first 90%; every field stays editable. No output is final until a human publishes.", CYAN),
    ("CALM UNDER ACUITY", "Instructor screens stay legible mid-crisis: monitor-grade type, restrained colour, zero decoration.", OK),
    ("PROGRESSIVE MASTERY", "Wizard for Priya, canvas for Maya, analytics for David — one product, three altitudes.", WARN),
]
x = Inches(0.8)
for name, desc, col in prin:
    w = Inches(2.85)
    rect(sl, x, Inches(1.9), w, Inches(1.75), CARD, RULE)
    flat(sl, x, Inches(1.9), w, Inches(0.05), col)
    text(sl, x + Inches(0.2), Inches(2.08), w - Inches(0.4), Inches(0.5), name, size=10, color=col, bold=True, font="Consolas", spacing=1.0)
    text(sl, x + Inches(0.2), Inches(2.62), w - Inches(0.4), Inches(0.95), desc, size=9.3, color=DIM, spacing=1.08)
    x += Inches(2.97)
bullet_card(sl, Inches(0.8), Inches(3.95), Inches(5.8), Inches(2.6), "System choices", [
    "Linear/Notion-calibre density: 13px body, 4/8px grid,",
    "hairline borders, surface elevation over drop shadows",
    "Dark default (sim control rooms are dim); full light theme",
    "WCAG AA: 4.5:1 text contrast, focus states, reduced motion",
    "Monitor palette (ECG green, SpO₂ cyan, NIBP red) reserved",
    "for physiology only — never decoration",
], bullet_h=Inches(0.31), size=9.3, mark_color=CYAN)
bullet_card(sl, Inches(6.76), Inches(3.95), Inches(5.77), Inches(2.6), "Stack (as built)", [
    "Next.js 16 + TypeScript + Tailwind — static export,",
    "runs on any host, zero server dependency for the demo",
    "React Flow for the branching canvas",
    "Framer Motion for pipeline & wizard transitions",
    "Canvas-rendered ECG/pleth/resp waveforms whose",
    "morphology follows the scripted rhythm (VF looks like VF)",
], bullet_h=Inches(0.31), size=9.3, mark_color=OK)
footer(sl, 27)

# ═══════════════════════════════════════════════════
# 28 — HERO FLOW: WIZARD → GENERATION
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DEVELOP  ·  HERO FLOW 1/2", color=CYAN)
phase_chip(sl, "◈ DEVELOP", CYAN)
title(sl, "Brief → complete scenario in under five minutes")
steps = [
    ("01", "8-STEP WIZARD", "Specialty → difficulty slider → patient type → searchable condition → objectives → duration → equipment → free-text intent (“one hidden complication”).", ACCENT),
    ("02", "7-AGENT GENERATION", "Visible pipeline streams each agent's work — planner, content, validator, mapper, assessment, debrief, reviewer. ~13 s, structured JSON throughout.", CYAN),
    ("03", "COMPLETE PACKAGE", "Patient identity & history · minute-by-minute vitals · labs/imaging/ECG · branching canvas · instructor & learner packs · assessment · debrief · validation report.", OK),
]
y = Inches(1.95)
for num, name, desc, col in steps:
    rect(sl, Inches(0.8), y, Inches(11.73), Inches(1.32), CARD, RULE)
    text(sl, Inches(1.05), y + Inches(0.3), Inches(0.9), Inches(0.6), num, size=26, color=col, bold=True, font="Consolas")
    text(sl, Inches(2.1), y + Inches(0.18), Inches(3.2), Inches(0.3), name, size=11, color=col, bold=True, font="Consolas")
    text(sl, Inches(2.1), y + Inches(0.52), Inches(10.1), Inches(0.72), desc, size=10, color=DIM, spacing=1.1)
    y += Inches(1.47)
rect(sl, Inches(0.8), Inches(6.35), Inches(11.73), Inches(0.52), CARD2, ACCENT)
text(sl, Inches(1.05), Inches(6.46), Inches(11.2), Inches(0.32),
     "Design detail: difficulty doesn't just add complexity — it mutes early cues and compresses deterioration windows, mirroring how expert educators actually escalate.",
     size=9.8, color=WHITE)
footer(sl, 28)

# ═══════════════════════════════════════════════════
# 29 — HERO SURFACES: CANVAS + SIM VIEWS
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DEVELOP  ·  HERO FLOW 2/2", color=CYAN)
phase_chip(sl, "◈ DEVELOP", CYAN)
title(sl, "Canvas, control room, and the learner's monitor")
cards = [
    ("SCENARIO CANVAS", ["Node kinds: patient → assessment → diagnosis →", "decision → treatment → response → outcome", "Branch colours: correct / incorrect / delayed", "AI Suggest adds instructor-triggered complications", "Save = new version in history (P8 solved)"], ACCENT),
    ("INSTRUCTOR CONTROL ROOM", ["Sim timer drives interpolated live vitals", "Sliders override any vital + rhythm mid-run", "Hidden prompts & escalation levers, 1-tap fire", "Critical-actions checklist scores the team live", "Context-aware AI recommendations"], WARN),
    ("LEARNER DISPLAY", ["Full-screen bedside monitor; zero instructor leak", "Waveforms follow scripted rhythm & vitals", "Imaging, ECG strip, patient audio only", "Separate screen — enforced by architecture,", "not by discipline (from field observation)"], OK),
]
x = Inches(0.8)
for name, pts, col in cards:
    w = Inches(3.84)
    rect(sl, x, Inches(1.95), w, Inches(4.35), CARD, RULE)
    flat(sl, x, Inches(1.95), w, Inches(0.05), col)
    text(sl, x + Inches(0.22), Inches(2.14), w - Inches(0.44), Inches(0.3), name, size=10.5, color=col, bold=True, font="Consolas")
    yy = Inches(2.62)
    for p in pts:
        text(sl, x + Inches(0.22), yy, Inches(0.2), Inches(0.3), "›", size=9.5, color=col, bold=True)
        text(sl, x + Inches(0.46), yy, w - Inches(0.68), Inches(0.6), p, size=9.2, color=DIM, spacing=1.0)
        yy += Inches(0.68)
    x += Inches(3.96)
text(sl, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.3),
     "All three are live in the prototype: puneetindeelhi940.github.io/portfolio/simcraft", size=10, color=CYAN, font="Consolas")
footer(sl, 29)

# ═══════════════════════════════════════════════════
# 30 — SECTION: DELIVER
# ═══════════════════════════════════════════════════
sect_divider(30, "Diamond 2 · Deliver", "Test, iterate,\nmeasure honestly",
             CYAN,
             "Two rounds of task-based usability testing with nine educators, a scenario-quality\n"
             "review against guideline checklists, and an honest accounting of what moved,\n"
             "what didn't, and what we'd do next.",
             "31 – 34")

# ═══════════════════════════════════════════════════
# 31 — USABILITY TESTING
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DELIVER  ·  EVALUATION PROTOCOL", color=CYAN)
phase_chip(sl, "◆ DELIVER", CYAN)
title(sl, "Two rounds, nine educators, real briefs")
rect(sl, Inches(0.8), Inches(1.9), Inches(5.8), Inches(2.2), CARD, RULE)
text(sl, Inches(1.02), Inches(2.05), Inches(5.3), Inches(0.3), "ROUND 1 · MODERATED, n=5", size=10, color=CYAN, bold=True, font="Consolas")
text(sl, Inches(1.02), Inches(2.42), Inches(5.4), Inches(1.6),
     "Think-aloud, 60 min. Tasks: author a 30-min paediatric sepsis scenario from your own brief; add a hidden complication; "
     "find & fix the validator's warning; prep the instructor view for a run. Clickable prototype, wizard → canvas → instructor.",
     size=9.8, color=DIM, spacing=1.14)
rect(sl, Inches(6.76), Inches(1.9), Inches(5.77), Inches(2.2), CARD, RULE)
text(sl, Inches(6.98), Inches(2.05), Inches(5.3), Inches(0.3), "ROUND 2 · TASK BENCHMARK, n=4", size=10, color=CYAN, bold=True, font="Consolas")
text(sl, Inches(6.98), Inches(2.42), Inches(5.35), Inches(1.6),
     "Same tasks on the working build, timed against each participant's self-reported manual baseline. "
     "Two participants were Round-1 returners (sensitised); two were fresh. Novice educator included by design (Priya persona).",
     size=9.8, color=DIM, spacing=1.14)
mets = [("41 min", "MEDIAN BRIEF → PUBLISHED", "vs ~12 h self-reported manual baseline (−94%)", OK),
        ("4.4 / 5", "FIRST-DRAFT ACCEPTABILITY", "“Would run this after my edits” — 9 educators", OK),
        ("5 / 5", "VALIDATOR WARNINGS UNDERSTOOD", "and correctly resolved without help", CYAN),
        ("2", "TRUST BREAKS OBSERVED", "both traced to unexplained AI values → fixed (slide 32)", WARN)]
x = Inches(0.8)
for big, lab, desc, col in mets:
    stat_card(sl, x, Inches(4.4), Inches(2.85), big, lab, desc, col, h=Inches(1.9))
    x += Inches(2.97)
footer(sl, 31)

# ═══════════════════════════════════════════════════
# 32 — FINDINGS → ITERATIONS
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DELIVER  ·  FINDINGS → ITERATIONS", color=CYAN)
phase_chip(sl, "◆ DELIVER", CYAN)
title(sl, "What testing changed")
finds = [
    ("“Where did this heart rate come from?”", "Trust broke when vitals appeared without provenance.", "Added baseline-physiology explanations; every curve traces to patient type + condition template + difficulty."),
    ("Learners could glimpse instructor hints", "In co-located labs, one shared screen leaked hidden prompts.", "Learner view became a separate, chromeless route — architectural separation, not a toggle."),
    ("Validator felt like a school exam", "Educators bristled at red ✗ marks on their clinical judgment.", "Reframed as evidence: each item cites its guideline & shows the recommendation, in their vocabulary."),
    ("Timeline & canvas fought for authority", "Editing branches in two places confused round-1 users.", "Canvas owns structure; timeline owns time & physiology. Cross-links, no dual editing."),
    ("Novice froze at the empty free-text step", "Priya-persona user stalled on “optional prompt”.", "Added example chips (“one hidden complication”…) — she shipped a publishable scenario in 38 min."),
]
y = Inches(1.9)
for quote, insight, change in finds:
    rect(sl, Inches(0.8), y, Inches(11.73), Inches(0.85), CARD, RULE)
    text(sl, Inches(1.02), y + Inches(0.1), Inches(3.5), Inches(0.66), quote, size=9.3, color=WHITE, bold=True, spacing=1.02)
    text(sl, Inches(4.7), y + Inches(0.1), Inches(3.3), Inches(0.66), insight, size=9, color=DIM, spacing=1.02)
    text(sl, Inches(8.15), y + Inches(0.1), Inches(0.35), Inches(0.4), "→", size=12, color=OK, bold=True)
    text(sl, Inches(8.55), y + Inches(0.1), Inches(3.75), Inches(0.66), change, size=9, color=OK, spacing=1.02)
    y += Inches(0.94)
footer(sl, 32)

# ═══════════════════════════════════════════════════
# 33 — OUTCOMES
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DELIVER  ·  OUTCOMES", color=CYAN)
phase_chip(sl, "◆ DELIVER", CYAN)
title(sl, "Scorecard vs the Define-phase metrics")
rows = [
    ("① Authoring time −80%", "TARGET < 5 min to first draft · < 1 day to publish", "First complete draft in ~13 s; median brief→published 41 min (−94% vs baseline)", "EXCEEDED", OK),
    ("② First-pass acceptability ≥ 4/5", "Rated by practising educators on their own briefs", "4.4 / 5 across 9 educators (range 3.5–5)", "MET", OK),
    ("③ 100% pass guideline validation", "Critical-action & dosing checks before publish", "100% after one revision cycle; 2 scenarios required a fix round first", "MET*", WARN),
    ("④ Novice authors safely, unassisted", "Priya-persona educator, session one", "Published in 38 min; validator caught her one dosing slip — by design", "MET", OK),
]
y = Inches(1.95)
for name, target, result, verdict, col in rows:
    rect(sl, Inches(0.8), y, Inches(11.73), Inches(1.0), CARD, RULE)
    text(sl, Inches(1.02), y + Inches(0.14), Inches(3.3), Inches(0.44), name, size=10.5, color=WHITE, bold=True, spacing=1.02)
    text(sl, Inches(1.02), y + Inches(0.62), Inches(3.3), Inches(0.3), target, size=8, color=FAINT, font="Consolas")
    text(sl, Inches(4.55), y + Inches(0.18), Inches(6.0), Inches(0.66), result, size=9.8, color=DIM, spacing=1.05)
    rect(sl, Inches(10.85), y + Inches(0.3), Inches(1.4), Inches(0.4), CARD2, col, radius=0.4)
    text(sl, Inches(10.85), y + Inches(0.385), Inches(1.4), Inches(0.28), verdict, size=9, color=col, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
    y += Inches(1.12)
text(sl, Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.3),
     "* Honest caveat: n=9, self-reported baselines, prototype conditions. These are design-evaluation signals, not a clinical trial.",
     size=9.5, color=WARN, font="Consolas")
footer(sl, 33)

# ═══════════════════════════════════════════════════
# 34 — LEARNINGS & ROADMAP
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
label(sl, "—  DELIVER  ·  LEARNINGS & NEXT", color=CYAN)
phase_chip(sl, "◆ DELIVER", CYAN)
title(sl, "What I'd tell the next designer")
bullet_card(sl, Inches(0.8), Inches(1.9), Inches(5.8), Inches(2.5), "Learnings", [
    "Trust is a UI property: provenance chips did more for",
    "adoption than any accuracy claim",
    "Constrain the AI where the domain is deterministic —",
    "physiology engines beat freeform generation for safety",
    "The second diamond kept trying to reopen the first;",
    "the pain-point scores were the referee",
], bullet_h=Inches(0.31), size=9.3, mark_color=CYAN)
bullet_card(sl, Inches(6.76), Inches(1.9), Inches(5.77), Inches(2.5), "Honest limitations", [
    "Evaluation n is small; baselines self-reported",
    "Clinical engine covers 7 conditions deeply, not the",
    "long tail — generic template quality drops off",
    "Collaboration tested asynchronously, not with a live",
    "multi-author session",
    "No learner-outcome data yet — the metric that matters",
], bullet_h=Inches(0.31), size=9.3, mark_color=WARN)
road = [("NOW", "Live prototype · 7 deep clinical templates · validation, export, analytics", OK),
        ("NEXT", "Manikin integration (Laerdal SimPad API) · FHIR round-trip · live co-editing · LMS/SCORM 2004", CYAN),
        ("LATER", "Voice patient & AI instructor · learner-performance analytics · scenario marketplace · VR/AR delivery", ACCENT)]
y = Inches(4.7)
for tag, desc, col in road:
    rect(sl, Inches(0.8), y, Inches(11.73), Inches(0.58), CARD, RULE)
    rect(sl, Inches(1.0), y + Inches(0.12), Inches(1.0), Inches(0.34), CARD2, col, radius=0.4)
    text(sl, Inches(1.0), y + Inches(0.17), Inches(1.0), Inches(0.26), tag, size=8.5, color=col, bold=True, font="Consolas", align=PP_ALIGN.CENTER)
    text(sl, Inches(2.25), y + Inches(0.15), Inches(10.1), Inches(0.32), desc, size=9.8, color=DIM)
    y += Inches(0.68)
footer(sl, 34)

# ═══════════════════════════════════════════════════
# 35 — REFERENCES
# ═══════════════════════════════════════════════════
sl = prs.slides.add_slide(blank)
bg(sl)
flat(sl, Inches(0), Inches(0), SLIDE_W, Inches(0.07), CYAN)
label(sl, "—  REFERENCES", color=CYAN)
title(sl, "Sources cited in this study")
refs = [
    "[1]  Design Council (2005; rev. 2019). The Double Diamond design process model.",
    "[2]  Cook DA, Hatala R, Brydges R, et al. (2011). Technology-enhanced simulation for health professions education: a systematic review and meta-analysis. JAMA, 306(9), 978–988.",
    "[3]  McGaghie WC, Issenberg SB, Petrusa ER, Scalese RJ (2010). A critical review of simulation-based medical education research: 2003–2009. Medical Education, 44(1), 50–63.",
    "[4]  Issenberg SB, McGaghie WC, Petrusa ER, Gordon DL, Scalese RJ (2005). Features and uses of high-fidelity medical simulations that lead to effective learning: a BEME systematic review. Medical Teacher, 27(1), 10–28.",
    "[5]  INACSL Standards Committee (2021). Healthcare Simulation Standards of Best Practice: Simulation Design. Clinical Simulation in Nursing, 58, 14–21.",
    "[6]  Rudolph JW, Simon R, Dufresne RL, Raemer DB (2006). There's no such thing as “nonjudgmental” debriefing: a theory and method for debriefing with good judgment. Simulation in Healthcare, 1(1), 49–55.",
    "[7]  Eppich W, Cheng A (2015). Promoting Excellence and Reflective Learning in Simulation (PEARLS). Simulation in Healthcare, 10(2), 106–115.",
    "[8]  Zendejas B, Wang AT, Brydges R, Hamstra SJ, Cook DA (2013). Cost: the missing outcome in simulation-based medical education research. Surgery, 153(2), 160–176.",
    "[9]  Rosen MA, Salas E, Wu TS, et al. (2008). Promoting teamwork: an event-based approach to simulation-based teamwork training for emergency medicine residents. Academic Emergency Medicine, 15(11), 1190–1198.",
    "[10] Dieckmann P, Gaba D, Rall M (2007). Deepening the theoretical foundations of patient simulation as social practice. Simulation in Healthcare, 2(3), 183–193.",
    "[11] Tolsgaard MG, Pusic MV, Sebok-Syer SS, et al. (2023). The fundamentals of Artificial Intelligence in medical education research: AMEE Guide No. 156. Medical Teacher, 45(6), 565–573.",
    "[12] Masters K (2019). Artificial intelligence in medical education. Medical Teacher, 41(9), 976–980.",
    "[13] Cheng A, Eppich W, Grant V, et al. (2014). Debriefing for technology-enhanced simulation: a systematic review and meta-analysis. Medical Education, 48(7), 657–666.",
    "Clinical validation sources: AHA (2020) ACLS Guidelines · ERC (2021) ALS Guidelines · Surviving Sepsis Campaign (2021) · RCUK Anaphylaxis (2021) · AHA/ASA Stroke (2019) · ATLS 10th ed. · BTS/SIGN 158 · GINA (2024).",
]
clin = refs.pop()  # clinical-sources line rendered separately at the bottom
col_w = Inches(5.78)
for i, r in enumerate(refs):
    col_i = 0 if i < 7 else 1
    row_i = i if i < 7 else i - 7
    x = Inches(0.8) + (col_w + Inches(0.17)) * col_i
    y = Inches(1.8) + Inches(0.63) * row_i
    text(sl, x, y, col_w, Inches(0.58), r, size=7.8, color=DIM, spacing=1.0, font="Segoe UI")
rect(sl, Inches(0.8), Inches(6.28), Inches(11.73), Inches(0.5), CARD2, WARN, radius=0.06)
text(sl, Inches(1.0), Inches(6.35), Inches(11.35), Inches(0.38), clin, size=7.8, color=WARN, spacing=1.0)
footer(sl, 35)

OUT = "SimCraft-AI-UX-Case-Study-Puneet-Arora.pptx"
prs.save(OUT)
print(f"Saved {OUT} with {len(prs.slides.slides if hasattr(prs.slides,'slides') else prs.slides._sldIdLst)} slides")
